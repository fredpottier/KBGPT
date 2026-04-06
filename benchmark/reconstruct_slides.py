#!/usr/bin/env python3
"""Reconstruit les slides PPTX en chunks autonomes pour le test de borne superieure.

Strategie :
1. Extraire le titre de la slide
2. Extraire tous les textes des shapes avec hierarchie (indentation)
3. Extraire les notes orateur
4. Combiner en un chunk unique prefixe par le contexte document/page
5. Tester avec Claude Sonnet et Qwen 14B
"""
import json
import os
import re
import time
import glob

# Les 13 questions PPTX du benchmark avec le slide index et le PPTX source
QUESTIONS = [
    {"qid": "T1_HUM_0040", "q": "Qu'est-ce que Joule dans le contexte SAP S/4HANA Cloud 2025 ?",
     "pptx": "022", "slide_keywords": ["Joule", "orchestrator"], "expected": "copilote,orchestrat,Business AI"},
    {"qid": "T1_HUM_0041", "q": "Quels partenaires AI sont listes dans l'ecosysteme SAP Business Data Cloud 2025 ?",
     "pptx": "022", "slide_keywords": ["ALEPH ALPHA", "Anthropic"], "expected": "Anthropic,Aleph Alpha,Cohere"},
    {"qid": "T1_HUM_0042", "q": "Que permet la fonctionnalite Create Purchase Requisitions de Joule ?",
     "pptx": "022", "slide_keywords": ["Purchase Requisition", "Joule"], "expected": "purchase requisition,material,service"},
    {"qid": "T1_HUM_0043", "q": "Qu'est-ce que le Production Planning Optimizer (PPO) ?",
     "pptx": "025", "slide_keywords": ["Production Planning Optimizer"], "expected": "PPO,constraint,linear programming"},
    {"qid": "T1_HUM_0044", "q": "Combien de pays de l'UE couverts par Intrastat dans S/4HANA ?",
     "pptx": "018", "slide_keywords": ["Intrastat", "member"], "expected": "27,member states"},
    {"qid": "T1_HUM_0045", "q": "Validite d'un permis de travail de niveau 1 dans Oil and Gas ?",
     "pptx": "025", "slide_keywords": ["Work Permit", "level 1"], "expected": "level 1,shift,permit"},
    {"qid": "T1_HUM_0052", "q": "Avantage de l'EWM embarque dans S/4HANA ?",
     "pptx": "018", "slide_keywords": ["EWM", "embedded", "eliminat"], "expected": "embedded,duplication,synchronization"},
    {"qid": "T1_HUM_0053", "q": "Combien de Lines of Business couvertes par S/4HANA 1809 ?",
     "pptx": "018", "slide_keywords": ["Lines of Business"], "expected": "10,Lines of Business"},
    {"qid": "T1_HUM_0063", "q": "Comment S/4HANA utilise la recherche fuzzy HANA pour la classification ?",
     "pptx": "025", "slide_keywords": ["fuzzy", "tariff"], "expected": "fuzzy,tariff,classification"},
    {"qid": "T1_HUM_0064", "q": "Les 12 steps du WIP management avec MES integration ?",
     "pptx": "025", "slide_keywords": ["WIP", "MES"], "expected": "MES,staging,warehouse"},
    {"qid": "T1_HUM_0070", "q": "SAP product ID pour S/4HANA Reinsurance for Assumed Risk ?",
     "pptx": "024", "slide_keywords": ["7019930", "Reinsurance"], "expected": "7019930,Reinsurance,Assumed Risk"},
    {"qid": "T1_HUM_0075", "q": "Comment fonctionne l'Advanced ATP (aATP) dans S/4HANA 1809 ?",
     "pptx": "018", "slide_keywords": ["Advanced ATP", "aATP"], "expected": "aATP,backorder,delivery"},
    {"qid": "T1_HUM_0081", "q": "Modeles de licensing pour S/4HANA Enterprise Management ?",
     "pptx": "024", "slide_keywords": ["Enterprise Management", "licensing"], "expected": "Enterprise Management,Enhanced LoB,Industry"},
]


def find_pptx_files():
    """Trouve tous les PPTX dans data/."""
    pptx_files = {}
    for root, dirs, files in os.walk("data"):
        for f in files:
            if f.endswith(".pptx"):
                path = os.path.join(root, f)
                pptx_files[f[:3]] = path  # Index par les 3 premiers chars (ex: "022")
                # Aussi indexer par prefixes plus longs
                for prefix in ["022", "018", "024", "025", "023", "007"]:
                    if f.startswith(prefix):
                        pptx_files[prefix] = path
    return pptx_files


def extract_slide_chunk(pptx_path, slide_keywords, doc_name=""):
    """Reconstruit le chunk d'une slide a partir des keywords."""
    from pptx import Presentation

    prs = Presentation(pptx_path)

    # Trouver la slide qui contient les keywords
    best_slide_idx = -1
    best_score = 0

    for idx, slide in enumerate(prs.slides):
        # Collecter tout le texte de la slide
        all_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text += " " + shape.text_frame.text

        # Scorer par nombre de keywords trouves
        score = sum(1 for kw in slide_keywords if kw.lower() in all_text.lower())
        if score > best_score:
            best_score = score
            best_slide_idx = idx

    if best_slide_idx < 0:
        return "", f"No slide found for keywords {slide_keywords}"

    slide = prs.slides[best_slide_idx]

    # Reconstruire le chunk de la slide
    parts = []

    # 1. Titre
    title = ""
    if slide.shapes.title:
        title = slide.shapes.title.text.strip()
        parts.append(f"## {title}")

    # 2. Contenu des shapes (avec hierarchie)
    content_lines = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text or text == title:
                    continue
                level = para.level
                if level > 0:
                    prefix = "  " * level + "- "
                else:
                    prefix = "- " if len(text) < 200 else ""
                content_lines.append(f"{prefix}{text}")

    if content_lines:
        parts.append("\n".join(content_lines))

    # 3. Notes orateur
    notes_text = ""
    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
        notes_text = slide.notes_slide.notes_text_frame.text.strip()
        if notes_text and len(notes_text) > 10:
            parts.append(f"\n--- Speaker Notes ---\n{notes_text}")

    # 4. Prefixe contextuel
    slide_content = "\n\n".join(parts)
    chunk = f"[Document: {doc_name} | Slide {best_slide_idx + 1}]\n\n{slide_content}"

    meta = f"slide {best_slide_idx + 1}, {len(slide_content)} chars, notes={'yes' if notes_text else 'no'}"
    return chunk[:3000], meta


def test_answer(answer, expected_str):
    terms = [t.strip().lower() for t in expected_str.split(",") if len(t.strip()) > 2]
    found = sum(1 for t in terms if t in answer.lower())
    return found / max(len(terms), 1), found, len(terms)


PROMPT = (
    "You are a precise assistant. Answer the question using ONLY the provided source. "
    "Be specific: include names, numbers, values, technical terms. "
    "If the information is partially available, answer with what you have. "
    "If the answer is not in the source at all, say 'information not available'. "
    "Answer in the SAME LANGUAGE as the question."
)


def main():
    pptx_files = find_pptx_files()
    print(f"PPTX files found: {list(pptx_files.keys())}")

    # Reconstruire les slides
    print("\nReconstruction des slides...")
    reconstructed = []
    for q in QUESTIONS:
        pptx_path = pptx_files.get(q["pptx"])
        if not pptx_path:
            print(f"  {q['qid']}: PPTX {q['pptx']} not found — will use cache fallback")
            reconstructed.append({"qid": q["qid"], "question": q["q"], "expected": q["expected"], "chunk": "", "meta": "not found"})
            continue

        doc_name = os.path.basename(pptx_path).replace(".pptx", "")
        chunk, meta = extract_slide_chunk(pptx_path, q["slide_keywords"], doc_name)
        reconstructed.append({"qid": q["qid"], "question": q["q"], "expected": q["expected"], "chunk": chunk, "meta": meta})
        print(f"  {q['qid']}: {meta}")

    # Fallback : pour les PPTX non trouves, utiliser le cache full_text
    for r in reconstructed:
        if not r["chunk"]:
            q = next(q for q in QUESTIONS if q["qid"] == r["qid"])
            # Chercher dans les caches
            for cache_file in glob.glob("data/extraction_cache/*.v5cache.json"):
                try:
                    data = json.load(open(cache_file, encoding="utf-8"))
                    ft = data.get("extraction", {}).get("full_text", "")
                    for kw in q["slide_keywords"]:
                        idx = ft.lower().find(kw.lower())
                        if idx >= 0:
                            pages = re.split(r'\[PAGE (\d+)\]', ft)
                            for i in range(1, len(pages), 2):
                                content = pages[i + 1] if i + 1 < len(pages) else ""
                                if kw.lower() in content.lower():
                                    clean = re.sub(r'\[PARAGRAPH\]', '\n', content)
                                    clean = re.sub(r'\[TITLE[^\]]*\]', '\n## ', clean)
                                    clean = re.sub(r'\[e\d+\|[^\]]*\]\s*"[^"]*"', '', clean)
                                    clean = re.sub(r'\[/?[A-Z_]+[^\]]*\]', '', clean)
                                    clean = re.sub(r'\n{3,}', '\n\n', clean).strip()
                                    # Window autour du keyword
                                    kidx = clean.lower().find(kw.lower())
                                    if kidx >= 0:
                                        start = max(0, kidx - 600)
                                        end = min(len(clean), kidx + 900)
                                        passage = clean[start:end].strip()
                                        r["chunk"] = f"[Document: cache | Page {pages[i]}]\n\n{passage}"[:2500]
                                        r["meta"] = f"cache page {pages[i]}, {len(passage)} chars"
                                        break
                            if r["chunk"]:
                                break
                except Exception:
                    pass
                if r["chunk"]:
                    break
            if r["chunk"]:
                print(f"  {r['qid']}: FALLBACK {r['meta']}")

    # Test Claude Sonnet
    print("\n=== CLAUDE SONNET ===")
    try:
        import anthropic
        client = anthropic.Anthropic()
        claude_results = []
        for r in reconstructed:
            if not r["chunk"]:
                claude_results.append({"qid": r["qid"], "status": "EMPTY", "score": 0})
                print(f"  [EMPTY] {r['qid']}: no chunk available")
                continue
            resp = client.messages.create(
                model="claude-sonnet-4-20250514", max_tokens=400, system=PROMPT,
                messages=[{"role": "user", "content": f"Source:\n{r['chunk']}\n\nQuestion: {r['question']}"}],
                temperature=0,
            )
            answer = resp.content[0].text if resp.content else ""
            score, found, total = test_answer(answer, r["expected"])
            status = "OK" if score >= 0.5 else "FAIL"
            claude_results.append({"qid": r["qid"], "status": status, "score": score, "answer": answer[:200]})
            preview = answer[:120].encode("ascii", "replace").decode().replace("\n", " ")
            print(f"  [{status}] {r['qid']}: {found}/{total} — {preview}")
            time.sleep(0.5)
        claude_ok = sum(1 for r in claude_results if r["status"] == "OK")
        print(f"\nClaude: {claude_ok}/13 ({100 * claude_ok / 13:.0f}%)")
    except Exception as e:
        print(f"Claude error: {e}")
        claude_ok = 0

    # Test Qwen 14B
    print("\n=== QWEN 14B ===")
    try:
        from openai import OpenAI
        qwen_client = OpenAI(api_key="EMPTY", base_url=os.environ.get("VLLM_URL", "http://18.194.28.167:8000") + "/v1")
        qwen_results = []
        for r in reconstructed:
            if not r["chunk"]:
                qwen_results.append({"qid": r["qid"], "status": "EMPTY", "score": 0})
                print(f"  [EMPTY] {r['qid']}: no chunk available")
                continue
            resp = qwen_client.chat.completions.create(
                model="Qwen/Qwen2.5-14B-Instruct-AWQ", max_tokens=400,
                messages=[{"role": "system", "content": PROMPT},
                          {"role": "user", "content": f"Source:\n{r['chunk']}\n\nQuestion: {r['question']}"}],
                temperature=0,
            )
            answer = resp.choices[0].message.content or ""
            score, found, total = test_answer(answer, r["expected"])
            status = "OK" if score >= 0.5 else "FAIL"
            qwen_results.append({"qid": r["qid"], "status": status, "score": score, "answer": answer[:200]})
            preview = answer[:120].encode("ascii", "replace").decode().replace("\n", " ")
            print(f"  [{status}] {r['qid']}: {found}/{total} — {preview}")
            time.sleep(0.3)
        qwen_ok = sum(1 for r in qwen_results if r["status"] == "OK")
        print(f"\nQwen: {qwen_ok}/13 ({100 * qwen_ok / 13:.0f}%)")
    except Exception as e:
        print(f"Qwen error: {e}")
        qwen_ok = 0

    print(f"\n{'=' * 60}")
    print(f"BORNE SUPERIEURE V3 (reconstruction slide PPTX)")
    print(f"  Claude Sonnet:             {claude_ok}/13 ({100 * claude_ok / 13:.0f}%)")
    print(f"  Qwen 14B:                  {qwen_ok}/13 ({100 * qwen_ok / 13:.0f}%)")
    print(f"  V2 (window full_text):     Claude 7/13 (54%), Qwen 2/13 (15%)")
    print(f"  Actuel (chunks atomiques): 2/13 (15%)")
    print(f"{'=' * 60}")


if __name__ == "__main__":
    main()
