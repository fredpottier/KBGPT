"""
Phase 0 - Extraction de Texte Brut

Extrait le texte des documents PDF/PPTX pour le POC.
Reutilise l'infrastructure OSMOSIS existante.
"""

import os
from pathlib import Path
from typing import Dict, Optional, Tuple
from dataclasses import dataclass


@dataclass
class ExtractionResult:
    """Resultat d'extraction de texte"""
    doc_id: str
    title: str
    text: str
    page_count: int
    char_count: int
    source_path: str
    success: bool
    error: Optional[str] = None


class TextExtractor:
    """
    Extracteur de texte pour le POC.
    Supporte PDF et PPTX.
    """

    def __init__(self, use_ocr: bool = False):
        """
        Args:
            use_ocr: Utiliser OCR pour les PDFs scannes
        """
        self.use_ocr = use_ocr

    def extract_from_pdf(self, pdf_path: str) -> ExtractionResult:
        """
        Extrait le texte d'un PDF.

        Args:
            pdf_path: Chemin vers le fichier PDF

        Returns:
            ExtractionResult avec le texte extrait
        """
        path = Path(pdf_path)
        doc_id = path.stem

        try:
            # Essayer d'utiliser PyMuPDF (fitz) s'il est disponible
            text, page_count = self._extract_with_pymupdf(pdf_path)

            if not text.strip() and self.use_ocr:
                # Fallback OCR si le PDF est scanne
                text, page_count = self._extract_with_ocr(pdf_path)

            return ExtractionResult(
                doc_id=doc_id,
                title=path.stem,
                text=text,
                page_count=page_count,
                char_count=len(text),
                source_path=str(path.absolute()),
                success=True
            )

        except Exception as e:
            return ExtractionResult(
                doc_id=doc_id,
                title=path.stem,
                text="",
                page_count=0,
                char_count=0,
                source_path=str(path.absolute()),
                success=False,
                error=str(e)
            )

    def _extract_with_pymupdf(self, pdf_path: str) -> Tuple[str, int]:
        """Extraction avec PyMuPDF"""
        try:
            import fitz  # PyMuPDF
        except ImportError:
            raise ImportError("PyMuPDF (fitz) non installe. Installer avec: pip install PyMuPDF")

        doc = fitz.open(pdf_path)
        texts = []
        for page in doc:
            texts.append(page.get_text())
        doc.close()

        return '\n\n'.join(texts), len(texts)

    def _extract_with_ocr(self, pdf_path: str) -> Tuple[str, int]:
        """Extraction OCR avec pytesseract"""
        try:
            import fitz
            from PIL import Image
            import pytesseract
            import io
        except ImportError:
            raise ImportError("Dependances OCR manquantes. Installer: pip install pytesseract Pillow")

        doc = fitz.open(pdf_path)
        texts = []

        for page in doc:
            # Convertir en image
            pix = page.get_pixmap(dpi=150)
            img_data = pix.tobytes("png")
            img = Image.open(io.BytesIO(img_data))

            # OCR
            text = pytesseract.image_to_string(img, lang='fra+eng')
            texts.append(text)

        doc.close()
        return '\n\n'.join(texts), len(texts)

    def extract_from_pptx(self, pptx_path: str) -> ExtractionResult:
        """
        Extrait le texte d'un PPTX.

        Args:
            pptx_path: Chemin vers le fichier PPTX

        Returns:
            ExtractionResult avec le texte extrait
        """
        path = Path(pptx_path)
        doc_id = path.stem

        try:
            from pptx import Presentation
        except ImportError:
            return ExtractionResult(
                doc_id=doc_id,
                title=path.stem,
                text="",
                page_count=0,
                char_count=0,
                source_path=str(path.absolute()),
                success=False,
                error="python-pptx non installe"
            )

        try:
            prs = Presentation(pptx_path)
            texts = []

            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                texts.append('\n'.join(slide_text))

            full_text = '\n\n--- SLIDE ---\n\n'.join(texts)

            return ExtractionResult(
                doc_id=doc_id,
                title=path.stem,
                text=full_text,
                page_count=len(prs.slides),
                char_count=len(full_text),
                source_path=str(path.absolute()),
                success=True
            )

        except Exception as e:
            return ExtractionResult(
                doc_id=doc_id,
                title=path.stem,
                text="",
                page_count=0,
                char_count=0,
                source_path=str(path.absolute()),
                success=False,
                error=str(e)
            )

    def extract(self, file_path: str) -> ExtractionResult:
        """
        Extrait le texte d'un fichier (auto-detect format).

        Args:
            file_path: Chemin vers le fichier

        Returns:
            ExtractionResult
        """
        path = Path(file_path)
        ext = path.suffix.lower()

        if ext == '.pdf':
            return self.extract_from_pdf(file_path)
        elif ext in ['.pptx', '.ppt']:
            return self.extract_from_pptx(file_path)
        else:
            return ExtractionResult(
                doc_id=path.stem,
                title=path.stem,
                text="",
                page_count=0,
                char_count=0,
                source_path=str(path.absolute()),
                success=False,
                error=f"Format non supporte: {ext}"
            )

    def extract_batch(self, file_paths: list) -> Dict[str, ExtractionResult]:
        """
        Extrait le texte de plusieurs fichiers.

        Args:
            file_paths: Liste des chemins

        Returns:
            Dict[doc_id -> ExtractionResult]
        """
        results = {}
        for path in file_paths:
            result = self.extract(path)
            results[result.doc_id] = result
        return results
