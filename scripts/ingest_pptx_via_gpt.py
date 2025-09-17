# ingest_pptx_via_gpt.py ‚Äî version am√©lior√©e avec contexte global & thumbnails

import base64
import json
import os
import shutil
import subprocess
import time
import uuid
from datetime import datetime, timezone
from pathlib import Path
from typing import List, Dict, Any
from concurrent.futures import ThreadPoolExecutor
from utils.solution_normalizer import normalize_solution_name

from langdetect import detect, DetectorFactory, LangDetectException
from pdf2image import convert_from_path
from pptx import Presentation
from PIL import Image
from qdrant_client.models import PointStruct
from utils.shared_clients import (
    ensure_qdrant_collection,
    get_openai_client,
    get_qdrant_client,
    get_sentence_transformer,
)

from import_logging import setup_logging
from prompt_registry import load_prompts, select_prompt, render_prompt


# --- Initialisation des chemins et variables globales ---
def get_project_root() -> Path:
    return Path(__file__).resolve().parent.parent


PROJECT_ROOT = get_project_root()
DOCS_IN = PROJECT_ROOT / "docs_in"
DOCS_DONE = PROJECT_ROOT / "public_files" / "presentations"
SLIDES_PNG = PROJECT_ROOT / "public_files" / "slides"
THUMBNAILS_DIR = PROJECT_ROOT / "public_files" / "thumbnails"
STATUS_DIR = PROJECT_ROOT / "status"
LOGS_DIR = PROJECT_ROOT / "logs"
MODELS_DIR = PROJECT_ROOT / "models"

os.environ.setdefault("HF_HOME", str(MODELS_DIR))

QDRANT_COLLECTION = os.getenv("QDRANT_COLLECTION", "sap_kb")
GPT_MODEL = os.getenv("GPT_MODEL", "gpt-4o")
EMB_MODEL_NAME = os.getenv("EMB_MODEL_NAME", "intfloat/multilingual-e5-base")
MAX_WORKERS = int(os.getenv("MAX_WORKERS", "2"))

logger = setup_logging(LOGS_DIR, "ingest_debug.log")
DetectorFactory.seed = 0

PROMPT_REGISTRY = load_prompts(str(PROJECT_ROOT / "config" / "prompts.yaml"))

# --- Fonctions utilitaires ---


# Cr√©e les dossiers n√©cessaires au projet
def ensure_dirs():
    for d in [
        DOCS_IN,
        DOCS_DONE,
        SLIDES_PNG,
        THUMBNAILS_DIR,
        STATUS_DIR,
        LOGS_DIR,
        MODELS_DIR,
    ]:
        d.mkdir(parents=True, exist_ok=True)


# Ex√©cute une commande syst√®me avec timeout
def run_cmd(cmd, timeout=120):
    try:
        subprocess.run(cmd, check=True, timeout=timeout)
        return True
    except Exception as e:
        logger.error(f"Command failed ({e}): {' '.join(cmd)}")
    return False


# Encode une image en base64
def encode_image_base64(path: Path) -> str:
    return base64.b64encode(path.read_bytes()).decode("utf-8")


# Nettoie la r√©ponse GPT (retire les balises Markdown)
def clean_gpt_response(raw: str) -> str:
    import re

    s = (raw or "").strip()
    s = re.sub(r"^```(?:json)?\s*", "", s)
    s = re.sub(r"\s*```$", "", s)
    return s.strip()


# D√©tecte la langue d'un texte (ISO2)
def get_language_iso2(text: str) -> str:
    try:
        return detect(text)
    except LangDetectException:
        return "en"


# Embedding des textes via SentenceTransformer
def embed_texts(texts: List[str]) -> List[List[float]]:
    batched = [f"passage: {t}" for t in texts]
    embeddings = model.encode(
        batched, normalize_embeddings=True, convert_to_numpy=True
    )
    return embeddings.tolist()


# D√©coupe les slides en batchs selon le nombre de tokens
def chunk_slides_by_tokens(slides_data, max_tokens):
    batches = []
    current_batch = []
    current_tokens = 0
    for slide in slides_data:
        slide_text = (slide.get("text", "") + "\n" + slide.get("notes", "")).strip()
        slide_tokens = estimate_tokens(slide_text)
        if current_tokens + slide_tokens > max_tokens and current_batch:
            batches.append(current_batch)
            current_batch = []
            current_tokens = 0
        current_batch.append(slide)
        current_tokens += slide_tokens
    if current_batch:
        batches.append(current_batch)
    return batches


# Estime le nombre de tokens dans un texte
def estimate_tokens(text: str) -> int:
    return int(len(text.split()) / 0.75)


# Normalise une URL publique
def normalize_public_url(url: str) -> str:
    if not url:
        return ""
    u = url.strip().rstrip("/")
    if not (u.startswith("http://") or u.startswith("https://")):
        u = "https://" + u
    return u


# R√©sout le chemin vers LibreOffice/soffice
def resolve_soffice_path() -> str:
    cand = os.getenv("SOFFICE_PATH", "").strip()
    if cand and Path(cand).exists():
        return cand
    found = shutil.which("soffice") or shutil.which("libreoffice")
    return found or "/usr/bin/soffice"


PUBLIC_URL = normalize_public_url(os.getenv("PUBLIC_URL", "sapkb.ngrok.app"))
SOFFICE_PATH = resolve_soffice_path()

# --- Initialisation des clients et mod√®les ---
openai_client = get_openai_client()
qdrant_client = get_qdrant_client()
model = get_sentence_transformer(EMB_MODEL_NAME)
embedding_size = model.get_sentence_embedding_dimension()
if embedding_size is None:
    raise RuntimeError("SentenceTransformer returned no embedding dimension")
EMB_SIZE = int(embedding_size)
ensure_qdrant_collection(QDRANT_COLLECTION, EMB_SIZE)

MAX_TOKENS_THRESHOLD = 40000
MAX_PARTIAL_TOKENS = 8000
MAX_SUMMARY_TOKENS = 60000

# --- Fonctions principales du pipeline ---


# Convertit un PPTX en PDF via LibreOffice
# Retourne le chemin du PDF g√©n√©r√©
def convert_pptx_to_pdf(pptx_path: Path, output_dir: Path) -> Path:
    output_dir.mkdir(parents=True, exist_ok=True)
    logger.info(f"üìÑ Conversion PPTX‚ÜíPDF: {pptx_path.name}")
    ok = run_cmd(
        [
            SOFFICE_PATH,
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(output_dir),
            str(pptx_path),
        ],
        timeout=600,
    )
    pdf_path = output_dir / (pptx_path.stem + ".pdf")
    if not ok or not pdf_path.exists():
        raise RuntimeError("LibreOffice conversion failed or PDF missing")
    logger.debug(f"PDF path: {pdf_path} (exists={pdf_path.exists()})")
    return pdf_path


# Extrait le texte et les notes de chaque slide du PPTX
def extract_notes_and_text(pptx_path: Path) -> List[Dict[str, Any]]:
    logger.info(f"üìä Extraction texte+notes du PPTX: {pptx_path.name}")
    prs = Presentation(str(pptx_path))
    slides_data = []
    for i, slide in enumerate(prs.slides, start=1):
        notes = ""
        if getattr(slide, "has_notes_slide", False):
            notes_slide = getattr(slide, "notes_slide", None)
            if notes_slide and hasattr(notes_slide, "notes_text_frame"):
                tf = notes_slide.notes_text_frame
                if tf and hasattr(tf, "text"):
                    notes = (tf.text or "").strip()
        texts = []
        for shape in slide.shapes:
            txt = getattr(shape, "text", None)
            if isinstance(txt, str) and txt.strip():
                texts.append(txt.strip())
        slides_data.append(
            {
                "slide_index": i,
                "text": "\n".join(texts),
                "notes": notes,
            }
        )
    logger.debug(f"Slides parsed: {len(slides_data)}")
    return slides_data


# G√©n√®re une miniature pour une image de slide
def generate_thumbnail(image_path: Path) -> Path:
    img = Image.open(image_path)
    img.thumbnail((900, 900), Image.Resampling.LANCZOS)
    thumb_path = THUMBNAILS_DIR / image_path.name
    thumb_path.parent.mkdir(parents=True, exist_ok=True)
    img.save(thumb_path, "PNG")
    return thumb_path


# D√©coupe un texte en chunks avec chevauchement
def recursive_chunk(text: str, max_len=400, overlap_ratio=0.15) -> List[str]:
    tokens = text.split()
    step = int(max_len * (1 - overlap_ratio))
    chunks = []
    for i in range(0, len(tokens), step):
        chunk = tokens[i : i + max_len]
        chunks.append(" ".join(chunk))
        if i + max_len >= len(tokens):
            break
    return chunks


# R√©sume un deck PPTX trop volumineux en plusieurs passes GPT
def summarize_large_pptx(slides_data: List[Dict[str, Any]]) -> str:
    all_text = "\n\n".join(
        (slide.get("text", "") + "\n" + slide.get("notes", "")).strip()
        for slide in slides_data
        if slide.get("text", "") or slide.get("notes", "")
    )
    total_tokens = estimate_tokens(all_text)
    if total_tokens <= MAX_TOKENS_THRESHOLD:
        return all_text
    document_type = "generic"
    deck_prompt_id, deck_template = select_prompt(
        PROMPT_REGISTRY, "pptx", f"deck.{document_type}"
    )
    batches = chunk_slides_by_tokens(slides_data, MAX_PARTIAL_TOKENS)
    partial_summaries = []
    for batch in batches:
        batch_text = "\n\n".join(
            (slide.get("text", "") + "\n" + slide.get("notes", "")).strip()
            for slide in batch
            if slide.get("text", "") or slide.get("notes", "")
        )
        prompt = render_prompt(
            deck_template, summary_text=batch_text[:40000], source_name="partial"
        )
        try:
            response = openai_client.chat.completions.create(
                model=GPT_MODEL,
                messages=[
                    {
                        "role": "system",
                        "content": "You are a precise summarization assistant.",
                    },
                    {"role": "user", "content": prompt},
                ],
                temperature=0.1,
                max_tokens=MAX_PARTIAL_TOKENS,
            )
            raw = getattr(response.choices[0].message, "content", "") or ""
            summary = clean_gpt_response(raw)
            partial_summaries.append(summary)
        except Exception as e:
            logger.error(f"‚ùå Partial summary error: {e}")
            continue
    final_summary = "\n".join(partial_summaries)
    if estimate_tokens(final_summary) > MAX_SUMMARY_TOKENS:
        prompt = render_prompt(
            deck_template,
            summary_text=final_summary[: MAX_SUMMARY_TOKENS * 2],
            source_name="global",
        )
        try:
            response = openai_client.chat.completions.create(
                model=GPT_MODEL,
                messages=[
                    {
                        "role": "system",
                        "content": "You are a precise summarization assistant.",
                    },
                    {"role": "user", "content": prompt},
                ],
                temperature=0.1,
                max_tokens=MAX_SUMMARY_TOKENS,
            )
            raw = getattr(response.choices[0].message, "content", "") or ""
            final_summary = clean_gpt_response(raw)
        except Exception as e:
            logger.error(f"‚ùå Global summary reduction error: {e}")
            final_summary = final_summary[: MAX_SUMMARY_TOKENS * 100]
    return final_summary


# Analyse globale du deck pour extraire r√©sum√© et m√©tadonn√©es (document, solution)
def analyze_deck_summary(
    slides_data: List[Dict[str, Any]], source_name: str, document_type: str = "default"
) -> dict:
    logger.info(f"üîç GPT: analyse du deck via texte extrait ‚Äî {source_name}")
    summary_text = summarize_large_pptx(slides_data)
    doc_type = document_type or "generic"
    deck_prompt_id, deck_template = select_prompt(
        PROMPT_REGISTRY, "pptx", f"deck.{doc_type}"
    )
    prompt = render_prompt(
        deck_template, summary_text=summary_text, source_name=source_name
    )
    try:
        response = openai_client.chat.completions.create(
            model=GPT_MODEL,
            messages=[
                {
                    "role": "system",
                    "content": "You are a precise SAP document metadata extraction assistant.",
                },
                {
                    "role": "user",
                    "content": prompt,
                },
            ],
            temperature=0.1,
            max_tokens=2000,
        )
        raw = getattr(response.choices[0].message, "content", "") or ""
        cleaned = clean_gpt_response(raw)
        result = json.loads(cleaned) if cleaned else {}
        if not isinstance(result, dict):
            result = {}
        summary = result.get("summary", "")
        metadata = result.get("metadata", {})
        document_meta = metadata.get("document", {})
        solution_meta = metadata.get("solution", {})
        # --- Normalisation des solutions ---
        raw_main = solution_meta.get("main", "") or solution_meta.get(
            "main_solution", ""
        )
        sol_id, canon_name = normalize_solution_name(raw_main)
        solution_meta["id"] = sol_id
        solution_meta["main"] = canon_name or raw_main
        normalized_supporting = []
        for supp in solution_meta.get("supporting", []):
            sid, canon = normalize_solution_name(supp)
            normalized_supporting.append(canon or supp)
        solution_meta["supporting"] = list(set(normalized_supporting))
        normalized_mentioned = []
        for ment in solution_meta.get("mentioned", []):
            sid, canon = normalize_solution_name(ment)
            normalized_mentioned.append(canon or ment)
        solution_meta["mentioned"] = list(set(normalized_mentioned))
        logger.debug(
            f"Deck summary + metadata keys: {list(result.keys()) if result else 'n/a'}"
        )
        result["_prompt_meta"] = {
            "document_type": doc_type,
            "deck_prompt_id": deck_prompt_id,
            "prompts_version": PROMPT_REGISTRY.get("version", "unknown"),
        }
        return {
            "summary": summary,
            "document": document_meta,
            "solution": solution_meta,
            "_prompt_meta": result["_prompt_meta"],
        }
    except Exception as e:
        logger.error(f"‚ùå GPT metadata error: {e}")
        return {}


# Analyse d'un slide via GPT + image, retourne les chunks enrichis
def ask_gpt_slide_analysis(
    image_path,
    deck_summary,
    slide_index,
    source_name,
    text,
    notes,
    document_type="default",
    deck_prompt_id="unknown",
    retries=2,
):
    img_b64 = base64.b64encode(image_path.read_bytes()).decode("utf-8")
    doc_type = document_type or "generic"
    slide_prompt_id, slide_template = select_prompt(
        PROMPT_REGISTRY, "pptx", f"slide.{doc_type}"
    )
    prompt_text = render_prompt(
        slide_template,
        deck_summary=deck_summary,
        slide_index=slide_index,
        source_name=source_name,
        text=text,
        notes=notes,
    )
    msg = [
        {
            "role": "system",
            "content": "You analyze slides with visuals deeply and coherently.",
        },
        {
            "role": "user",
            "content": [
                {"type": "text", "text": prompt_text},
                {
                    "type": "image_url",
                    "image_url": {"url": f"data:image/png;base64,{img_b64}"},
                },
            ],
        },
    ]
    for attempt in range(retries):
        try:
            resp = openai_client.chat.completions.create(
                model=GPT_MODEL, messages=msg, temperature=0.2, max_tokens=1024
            )
            logger.debug(
                f"GPT response for slide {slide_index}: {resp.choices[0].message.content!r}"
            )
            raw_content = resp.choices[0].message.content
            cleaned_content = clean_gpt_response(raw_content or "")
            items = json.loads(cleaned_content)
            enriched = []
            for it in items:
                expl = it.get("full_explanation", "")
                meta = it.get("meta", {})
                if expl:
                    for seg in recursive_chunk(expl, max_len=400, overlap_ratio=0.15):
                        enriched.append(
                            {
                                "full_explanation": seg,
                                "meta": meta,
                                "prompt_meta": {
                                    "document_type": doc_type,
                                    "slide_prompt_id": slide_prompt_id,
                                    "prompts_version": PROMPT_REGISTRY.get(
                                        "version", "unknown"
                                    ),
                                },
                            }
                        )
            return enriched
        except Exception as e:
            logger.warning(f"Slide {slide_index} attempt {attempt} failed: {e}")
            time.sleep(2 * (attempt + 1))
    return []


# Ingestion des chunks dans Qdrant avec sch√©ma canonique
def ingest_chunks(chunks, doc_meta, file_uid, slide_index, deck_summary):
    valid = [ch for ch in chunks if ch.get("full_explanation", "").strip()]
    if not valid:
        logger.info(f"Slide {slide_index}: no valid chunks")
        return
    texts = [ch["full_explanation"] for ch in valid]
    embs = embed_texts(texts)
    points = []
    for ch, emb in zip(valid, embs):
        meta = ch.get("meta", {})
        payload = {
            "text": ch["full_explanation"].strip(),
            "language": get_language_iso2(ch["full_explanation"]),
            "ingested_at": datetime.now(timezone.utc).isoformat(),
            "document": {
                "source_name": f"{file_uid}.pptx",
                "source_type": "pptx",
                "source_file_url": f"{PUBLIC_URL}/static/presentations/{file_uid}.pptx",
                "slide_image_url": f"{PUBLIC_URL}/static/thumbnails/{file_uid}_slide_{slide_index}.png",
                "title": doc_meta.get("title", ""),
                "objective": doc_meta.get("objective", ""),
                "audience": doc_meta.get("audience", []),
                "source_date": doc_meta.get("source_date", ""),
            },
            "solution": {
                "main": doc_meta.get("main_solution", ""),
                "family": doc_meta.get("family", ""),
                "supporting": doc_meta.get("supporting_solutions", []),
                "mentioned": doc_meta.get("mentioned_solutions", []),
                "version": doc_meta.get("version", ""),
                "deployment_model": doc_meta.get("deployment_model", ""),
            },
            "chunk": {
                "scope": meta.get("scope", "solution-specific"),
                "slide_index": slide_index,
                "type": meta.get("type", ""),
                "level": meta.get("level", ""),
                "tags": meta.get("tags", []),
            },
            "deck_summary": deck_summary,
            "prompt_meta": ch.get("prompt_meta", {}),
        }
        points.append(PointStruct(id=str(uuid.uuid4()), vector=emb, payload=payload))
    qdrant_client.upsert(collection_name=QDRANT_COLLECTION, points=points)
    logger.info(f"Slide {slide_index}: ingested {len(points)} chunks")


# Fonction principale pour traiter un fichier PPTX
def process_pptx(pptx_path: Path, document_type: str = "default"):
    logger.info(f"start ingestion for {pptx_path.name}")
    ensure_dirs()
    pdf_path = convert_pptx_to_pdf(pptx_path, SLIDES_PNG)
    slides_data = extract_notes_and_text(pptx_path)
    deck_info = analyze_deck_summary(
        slides_data, pptx_path.name, document_type=document_type
    )
    summary = deck_info.get("summary", "")
    doc_meta = deck_info.get("document", {})
    solution_meta = deck_info.get("solution", {})
    deck_prompt_id = deck_info.get("_prompt_meta", {}).get("deck_prompt_id", "unknown")
    images = convert_from_path(str(pdf_path), output_folder=str(SLIDES_PNG))
    image_paths = {}
    for i, img in enumerate(images, start=1):
        img_path = THUMBNAILS_DIR / f"{pptx_path.stem}_slide_{i}.png"
        img.save(img_path, "PNG")
        image_paths[i] = img_path
        generate_thumbnail(img_path)
    tasks = []
    with ThreadPoolExecutor(max_workers=MAX_WORKERS) as ex:
        for slide in slides_data:
            idx = slide["slide_index"]
            text = slide["text"]
            notes = slide["notes"]
            if idx in image_paths:
                tasks.append(
                    (
                        idx,
                        ex.submit(
                            ask_gpt_slide_analysis,
                            image_paths[idx],
                            summary,
                            idx,
                            pptx_path.name,
                            text,
                            notes,
                            document_type,
                            deck_prompt_id,
                        ),
                    )
                )
    total = 0
    for idx, future in tasks:
        chunks = future.result() or []
        merged_meta = {**doc_meta, **solution_meta}
        ingest_chunks(chunks, merged_meta, pptx_path.stem, idx, summary)
        total += len(chunks)
    shutil.move(str(pptx_path), DOCS_DONE / f"{pptx_path.stem}.pptx")
    logger.info(f"Done {pptx_path.name} ‚Äî total chunks: {total}")


# Fusionne plusieurs dictionnaires de m√©tadonn√©es et normalise les solutions
def merge_metadata(meta_list: List[Dict[str, Any]]) -> Dict[str, Any]:
    merged = {
        "main_solution": "",
        "supporting_solutions": [],
        "mentioned_solutions": [],
        "document_type": "",
        "audience": [],
        "source_date": "",
        "language": "",
        "objective": "",
        "title": "",
        "version": "",
        "family": "",
        "deployment_model": "",
    }
    for meta in meta_list:
        for k in ["supporting_solutions", "mentioned_solutions", "audience"]:
            merged[k].extend(meta.get(k, []))
        for k in [
            "main_solution",
            "document_type",
            "source_date",
            "language",
            "objective",
            "title",
            "version",
            "family",
            "deployment_model",
        ]:
            if not merged[k] and meta.get(k):
                merged[k] = meta[k]
    for k in ["supporting_solutions", "mentioned_solutions", "audience"]:
        merged[k] = list(set(merged[k]))
    if merged.get("main_solution"):
        sol_id, canon = normalize_solution_name(merged["main_solution"])
        merged["main_solution_id"] = sol_id
        merged["main_solution"] = canon or merged["main_solution"]
    else:
        merged["main_solution_id"] = "UNMAPPED"
    normalized_supporting = []
    for supp in merged["supporting_solutions"]:
        sid, canon = normalize_solution_name(supp)
        normalized_supporting.append(canon or supp)
    merged["supporting_solutions"] = list(set(normalized_supporting))
    normalized_mentioned = []
    for ment in merged["mentioned_solutions"]:
        sid, canon = normalize_solution_name(ment)
        normalized_mentioned.append(canon or ment)
    merged["mentioned_solutions"] = list(set(normalized_mentioned))
    return merged


# Point d'entr√©e principal du script
def main():
    ensure_dirs()
    import argparse

    parser = argparse.ArgumentParser()
    parser.add_argument("pptx_path", type=str, help="Chemin du fichier PPTX √† ing√©rer")
    parser.add_argument(
        "--document-type",
        type=str,
        default=None,
        help="Type de document pour le choix des prompts",
    )
    args = parser.parse_args()
    pptx_path = Path(args.pptx_path)
    document_type = args.document_type
    if document_type is None:
        meta_path = pptx_path.with_suffix(".meta.json")
        if meta_path.exists():
            try:
                with open(meta_path, "r", encoding="utf-8") as f:
                    meta = json.load(f)
                document_type = meta.get("document_type", "default")
            except Exception as e:
                logger.warning(
                    f"Impossible de lire le document_type depuis {meta_path}: {e}"
                )
                document_type = "default"
        else:
            document_type = "default"
    if pptx_path.exists():
        process_pptx(pptx_path, document_type=document_type)
    else:
        logger.error(f"File not found: {pptx_path}")


if __name__ == "__main__":
    main()
