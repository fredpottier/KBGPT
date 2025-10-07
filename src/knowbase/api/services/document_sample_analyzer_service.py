"""
Service pour analyser des documents samples et suggérer entity types.

Phase 6 - Document Types Management
"""
import json
from typing import List, Dict, Optional
from fastapi import UploadFile
from pathlib import Path
import tempfile

from knowbase.common.llm_router import LLMRouter, TaskType
from knowbase.config.settings import get_settings
from knowbase.common.logging import setup_logging

settings = get_settings()
logger = setup_logging(settings.logs_dir, "document_sample_analyzer.log")


class DocumentSampleAnalyzerService:
    """Service pour analyser documents samples et suggérer entity types."""

    def __init__(self, llm_router: Optional[LLMRouter] = None, db_session=None):
        """
        Initialize service.

        Args:
            llm_router: Router LLM (optionnel, créé si None)
            db_session: Session DB pour récupérer entity types existants
        """
        self.llm_router = llm_router or LLMRouter()
        self.db_session = db_session

    async def analyze_document_sample(
        self,
        file: UploadFile,
        context_prompt: Optional[str] = None,
        model_preference: str = "claude-sonnet",
        tenant_id: str = "default"
    ) -> Dict:
        """
        Analyser un document sample pour suggérer entity types.

        Args:
            file: Fichier uploadé (PPTX ou PDF)
            context_prompt: Contexte additionnel pour guider le LLM
            model_preference: Modèle LLM à utiliser
            tenant_id: Tenant ID pour récupérer entity types existants

        Returns:
            Dict avec suggested_types, document_summary, pages_analyzed
        """
        logger.info(f"📄 Analyse document sample: {file.filename}")

        # Récupérer entity types existants approuvés
        existing_types = self._get_existing_entity_types(tenant_id)
        logger.info(f"📋 {len(existing_types)} entity types existants trouvés")

        # Extraire contenu du document
        content = await self._extract_document_content(file)

        if not content:
            logger.warning("⚠️ Aucun contenu extrait du document")
            return {
                "suggested_types": [],
                "document_summary": "Impossible d'extraire le contenu du document",
                "pages_analyzed": 0
            }

        # Construire prompt pour LLM
        prompt = self._build_analysis_prompt(content, context_prompt, existing_types)

        # Appeler LLM
        try:
            messages = [
                {
                    "role": "user",
                    "content": prompt
                }
            ]

            response = self.llm_router.complete(
                task_type=TaskType.METADATA_EXTRACTION,
                messages=messages,
                temperature=0.3,  # Température moyenne pour équilibre
                max_tokens=4000,
                model_preference=model_preference
            )

            # Parser réponse LLM
            result = self._parse_llm_response(response, existing_types)

            logger.info(
                f"✅ Analyse terminée: {len(result['suggested_types'])} types suggérés"
            )

            return result

        except Exception as e:
            logger.error(f"❌ Erreur analyse LLM: {e}")
            raise

    def _get_existing_entity_types(self, tenant_id: str = "default") -> List[str]:
        """
        Récupérer les entity types existants approuvés depuis la base.

        Args:
            tenant_id: Tenant ID

        Returns:
            Liste des noms de types approuvés
        """
        if not self.db_session:
            logger.warning("⚠️ Pas de session DB, impossible de récupérer entity types existants")
            return []

        try:
            from knowbase.db import DynamicEntityType

            # Requête pour récupérer types approuvés
            types = self.db_session.query(DynamicEntityType).filter(
                DynamicEntityType.tenant_id == tenant_id,
                DynamicEntityType.status == "approved"
            ).all()

            return [t.type_name for t in types]

        except Exception as e:
            logger.error(f"❌ Erreur récupération entity types existants: {e}")
            return []

    async def _extract_document_content(self, file: UploadFile) -> str:
        """
        Extraire contenu textuel du document.

        Args:
            file: Fichier uploadé

        Returns:
            Contenu textuel extrait
        """
        file_ext = Path(file.filename).suffix.lower()

        # Sauvegarder temporairement le fichier
        with tempfile.NamedTemporaryFile(delete=False, suffix=file_ext) as tmp_file:
            content = await file.read()
            tmp_file.write(content)
            tmp_path = Path(tmp_file.name)

        try:
            if file_ext == ".pdf":
                return await self._extract_pdf_content(tmp_path)
            elif file_ext in [".pptx", ".ppt"]:
                return await self._extract_pptx_content(tmp_path)
            else:
                logger.warning(f"⚠️ Format non supporté: {file_ext}")
                return ""
        finally:
            # Nettoyer fichier temporaire
            if tmp_path.exists():
                tmp_path.unlink()

    async def _extract_pdf_content(self, pdf_path: Path) -> str:
        """Extraire contenu d'un PDF."""
        try:
            from knowbase.ingestion.pipelines.pdf_pipeline import extract_text_from_pdf

            # Extraire texte (limiter à 10 premières pages pour analyse)
            text_parts = []
            pages = extract_text_from_pdf(str(pdf_path), max_pages=10)

            for page_num, page_text in pages:
                if page_text.strip():
                    text_parts.append(f"--- Page {page_num} ---\n{page_text}")

            return "\n\n".join(text_parts)

        except Exception as e:
            logger.error(f"❌ Erreur extraction PDF: {e}")
            return ""

    async def _extract_pptx_content(self, pptx_path: Path) -> str:
        """Extraire contenu d'un PPTX."""
        try:
            from pptx import Presentation

            prs = Presentation(str(pptx_path))
            text_parts = []

            # Limiter à 10 premières slides
            for i, slide in enumerate(prs.slides[:10], start=1):
                slide_text_parts = []

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text_parts.append(shape.text)

                if slide_text_parts:
                    text_parts.append(
                        f"--- Slide {i} ---\n" + "\n".join(slide_text_parts)
                    )

            return "\n\n".join(text_parts)

        except Exception as e:
            logger.error(f"❌ Erreur extraction PPTX: {e}")
            return ""

    def _build_analysis_prompt(
        self,
        document_content: str,
        context_prompt: Optional[str] = None,
        existing_types: Optional[List[str]] = None
    ) -> str:
        """
        Construire prompt LLM pour analyse.

        Args:
            document_content: Contenu du document
            context_prompt: Contexte additionnel
            existing_types: Liste des entity types déjà existants en base

        Returns:
            Prompt formaté
        """
        # Limiter taille du contenu (max 8000 chars)
        content_preview = document_content[:8000]
        if len(document_content) > 8000:
            content_preview += "\n\n[... contenu tronqué ...]"

        # Formater liste types existants
        existing_types_section = ""
        if existing_types:
            types_list = ", ".join(existing_types)
            existing_types_section = f"""
**TYPES D'ENTITÉS DÉJÀ EXISTANTS** (à privilégier si pertinents):
{types_list}

IMPORTANT: Si un type existe déjà dans cette liste, utilise-le tel quel (même nom exact).
Propose de NOUVEAUX types uniquement si aucun type existant ne convient.
"""

        context_section = ""
        if context_prompt:
            context_section = f"**CONTEXTE ADDITIONNEL**:\n{context_prompt}\n\n"

        prompt = f"""Tu es un expert en analyse de documents et modélisation d'entités.

**TÂCHE**: Analyser ce document et identifier les types d'entités pertinents pour l'extraction automatique.
{existing_types_section}
{context_section}**DOCUMENT À ANALYSER**:
{content_preview}

**INSTRUCTIONS**:
1. Analyse le document et identifie les types d'entités présents
2. **EN PRIORITÉ**: Vérifie si des types EXISTANTS (listés ci-dessus) sont pertinents
3. Propose de NOUVEAUX types uniquement si nécessaire
4. Pour chaque type identifié:
   - Donne un nom en MAJUSCULES (ex: SOLUTION, PRODUCT, INFRASTRUCTURE)
   - Estime une confidence (0.0 à 1.0) selon la fréquence et l'importance
   - Fournis 2-4 exemples concrets trouvés dans le document
   - Donne une brève description du type (1 phrase)
   - Indique "is_existing": true si le type est dans la liste existante, false sinon

5. **RÈGLES CRITIQUES**:
   - Réutilise les types existants EXACTEMENT (même nom, même casse)
   - Ne propose que des types **réellement présents** dans le document
   - Privilégie les types avec forte occurrence (confidence > 0.6)
   - Évite les types trop génériques (ex: "INFORMATION", "ELEMENT")
   - Maximum 15 types suggérés

6. **RÉSUMÉ**: Fournis aussi un résumé du document en 2-3 phrases

**FORMAT DE SORTIE JSON**:
{{
  "document_summary": "Ce document présente...",
  "suggested_types": [
    {{
      "name": "SOLUTION",
      "confidence": 0.95,
      "examples": ["SAP HANA", "SAP S/4HANA Cloud"],
      "description": "Solutions logicielles SAP mentionnées",
      "is_existing": true
    }},
    {{
      "name": "DEPLOYMENT_MODEL",
      "confidence": 0.85,
      "examples": ["Cloud Native", "Hybrid Deployment"],
      "description": "Modèles de déploiement mentionnés",
      "is_existing": false
    }}
  ]
}}

Retourne UNIQUEMENT le JSON, sans texte avant/après.
"""
        return prompt

    def _parse_llm_response(self, llm_response: str, existing_types: Optional[List[str]] = None) -> Dict:
        """
        Parser réponse LLM.

        Args:
            llm_response: Réponse brute LLM
            existing_types: Liste des types existants pour validation

        Returns:
            Dict avec suggested_types, document_summary, pages_analyzed
        """
        existing_types_set = set(existing_types) if existing_types else set()
        try:
            # Nettoyer réponse (supprimer markdown code blocks si présents)
            clean_response = llm_response.strip()
            if clean_response.startswith("```json"):
                clean_response = clean_response[7:]
            if clean_response.startswith("```"):
                clean_response = clean_response[3:]
            if clean_response.endswith("```"):
                clean_response = clean_response[:-3]

            clean_response = clean_response.strip()

            # Parser JSON
            data = json.loads(clean_response)

            # Valider structure
            suggested_types = data.get("suggested_types", [])
            document_summary = data.get("document_summary", "")

            # Valider chaque type suggéré
            validated_types = []
            for item in suggested_types:
                if not isinstance(item, dict):
                    continue

                # Champs requis
                if "name" not in item or "confidence" not in item:
                    continue

                # Normaliser nom (UPPERCASE)
                item["name"] = item["name"].upper().strip()

                # Valider confidence
                try:
                    item["confidence"] = float(item["confidence"])
                    if not (0.0 <= item["confidence"] <= 1.0):
                        item["confidence"] = 0.5
                except (TypeError, ValueError):
                    item["confidence"] = 0.5

                # Assurer que examples est une liste
                if "examples" not in item or not isinstance(item["examples"], list):
                    item["examples"] = []

                # Assurer que description existe
                if "description" not in item:
                    item["description"] = ""

                # Valider et forcer is_existing basé sur existing_types_set
                # (pour éviter que le LLM ne se trompe)
                if item["name"] in existing_types_set:
                    item["is_existing"] = True
                else:
                    item["is_existing"] = False

                validated_types.append(item)

            # Trier par confidence décroissant
            validated_types.sort(key=lambda x: x["confidence"], reverse=True)

            logger.info(f"✅ {len(validated_types)} types validés depuis réponse LLM")

            return {
                "suggested_types": validated_types,
                "document_summary": document_summary,
                "pages_analyzed": 10  # Max pages analysées
            }

        except json.JSONDecodeError as e:
            logger.error(f"❌ Erreur parsing JSON LLM: {e}")
            logger.error(f"Réponse LLM: {llm_response[:500]}")
            raise ValueError(f"LLM response is not valid JSON: {e}")

        except Exception as e:
            logger.error(f"❌ Erreur parsing réponse LLM: {e}")
            raise


__all__ = ["DocumentSampleAnalyzerService"]
