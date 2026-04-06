"""
Extracteur PPTX natif — bypass Docling pour les presentations PowerPoint.

Utilise python-pptx directement pour preserver :
- Hierarchie titre / sous-titre / bullets (placeholders + para.level)
- Speaker notes (verbatim auteur, mine d'or pour le retrieval)
- Tracking de sections (slides de titre detectees)
- Tables structurees

Produit des VisionUnit (compatibilite pipeline V2) et un
StructuralGraphBuildResult avec TypeAwareChunks = slides reconstruites.

Decision architecturale : Docling PPTX = SimplePipeline = wrapper python-pptx
degrade (bugs #1324 hierarchie, #1325 notes ignorees, #2551 titres absents).
Notre extracteur custom donne un controle total sur la qualite des chunks.
"""
from __future__ import annotations

import hashlib
import logging
import re
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple
from uuid import uuid4

from knowbase.extraction_v2.models.elements import BoundingBox, TextBlock, TableData
from knowbase.extraction_v2.models.vision_unit import VisionUnit
from knowbase.extraction_v2.models.extraction_result import (
    PageIndex,
    PageOutput,
    DocumentStructure,
)
from knowbase.structural.models import (
    ChunkKind,
    DocItem,
    DocItemType,
    DocumentVersion,
    PageContext,
    SectionInfo,
    TextOrigin,
    TypeAwareChunk,
)
from knowbase.structural.graph_builder import StructuralGraphBuildResult

logger = logging.getLogger(__name__)

# Seuil pour considerer une slide comme "titre de section"
MIN_CONTENT_CHARS = 50

# Dimensions PPTX par defaut (16:9 HD)
DEFAULT_SLIDE_WIDTH = 960.0
DEFAULT_SLIDE_HEIGHT = 540.0


@dataclass
class PptxExtractionResult:
    """Resultat complet de l'extraction PPTX."""
    units: List[VisionUnit]
    structural_graph: StructuralGraphBuildResult
    full_text: str
    structure: DocumentStructure
    page_index: List[PageIndex]
    doc_title: str


def _clean_whitespace(text: str) -> str:
    """Normalise espaces multiples et tabs."""
    text = text.replace("\t", " ")
    text = re.sub(r" {2,}", " ", text)
    return text.strip()


def _extract_title_and_subtitle(slide) -> Tuple[str, str]:
    """Extrait titre et sous-titre via les placeholders PPTX.

    PPTX utilise \\x0b (vertical tab) comme saut de ligne dans les titres.
    """
    title = ""
    subtitle = ""

    for shape in slide.placeholders:
        phf = shape.placeholder_format
        if phf is None:
            continue
        if phf.idx == 0 and shape.has_text_frame:
            raw = shape.text_frame.text.strip()
            if "\x0b" in raw:
                parts = raw.split("\x0b", 1)
                title = parts[0].strip()
                subtitle = parts[1].strip()
            else:
                title = raw
        elif phf.idx == 1 and shape.has_text_frame:
            candidate = shape.text_frame.text.strip()
            if candidate and len(candidate) < 200 and candidate.count("\n") <= 2:
                if not subtitle:
                    subtitle = candidate

    if not title and slide.shapes.title:
        raw = slide.shapes.title.text.strip()
        if "\x0b" in raw:
            parts = raw.split("\x0b", 1)
            title = parts[0].strip()
            if not subtitle:
                subtitle = parts[1].strip()
        else:
            title = raw

    return title, subtitle


def _get_title_shape_ids(slide) -> set:
    """Identifie les shape IDs des placeholders titre/sous-titre."""
    ids: set = set()
    for shape in slide.placeholders:
        phf = shape.placeholder_format
        if phf is not None and phf.idx in (0, 1):
            ids.add(shape.shape_id)
    return ids


def _extract_notes(slide) -> str:
    """Extrait les speaker notes."""
    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
        return slide.notes_slide.notes_text_frame.text.strip()
    return ""


def _extract_content(slide, title: str, subtitle: str) -> Tuple[List[TextBlock], List[TableData]]:
    """Extrait contenu (bullets + tables) en preservant la hierarchie."""
    skip_texts = {title, subtitle}
    title_shape_ids = _get_title_shape_ids(slide)
    blocks: List[TextBlock] = []
    tables: List[TableData] = []

    for shape in slide.shapes:
        if shape.shape_id in title_shape_ids:
            continue

        if shape.has_text_frame:
            shape_lines: List[Tuple[str, int]] = []  # (text, level)
            for para in shape.text_frame.paragraphs:
                text = _clean_whitespace(para.text)
                if not text or text in skip_texts:
                    continue
                shape_lines.append((text, para.level))

            # Filtrer les shapes qui ne contiennent que des labels courts
            if shape_lines:
                avg_len = sum(len(t) for t, _ in shape_lines) / len(shape_lines)
                if avg_len < 15 and len(shape_lines) > 2:
                    continue

            for text, level in shape_lines:
                block_type = "list_item" if level > 0 else "paragraph"
                blocks.append(TextBlock(
                    type=block_type,
                    text=text,
                    level=level,
                    block_id=f"block_{uuid4().hex[:8]}",
                ))

        if shape.has_table:
            table = shape.table
            rows_data = list(table.rows)
            if not rows_data:
                continue

            headers = [_clean_whitespace(cell.text) for cell in rows_data[0].cells]
            data_rows = []
            for row in rows_data[1:]:
                data_rows.append([_clean_whitespace(cell.text) for cell in row.cells])

            tables.append(TableData(
                table_id=f"table_{uuid4().hex[:8]}",
                headers=headers,
                cells=[headers] + data_rows,
                num_rows=len(data_rows) + 1,
                num_cols=len(headers),
            ))

    return blocks, tables


def _build_reconstructed_text(
    title: str,
    subtitle: str,
    blocks: List[TextBlock],
    tables: List[TableData],
    notes: str,
) -> str:
    """Construit le texte reconstruit d'une slide (sans prefixe contextuel)."""
    parts: List[str] = []

    if title:
        header = f"## {title}"
        if subtitle:
            header += f"\n### {subtitle}"
        parts.append(header)
    elif subtitle:
        parts.append(f"### {subtitle}")

    # Bullets avec hierarchie
    content_lines: List[str] = []
    for block in blocks:
        if block.level > 0:
            prefix = "  " * block.level + "- "
        else:
            prefix = "- " if len(block.text) < 300 else ""
        content_lines.append(f"{prefix}{block.text}")

    if content_lines:
        parts.append("\n".join(content_lines))

    # Tables en markdown
    for table in tables:
        table_lines = []
        table_lines.append("| " + " | ".join(table.headers) + " |")
        table_lines.append("| " + " | ".join(["---"] * len(table.headers)) + " |")
        # cells[0] = headers (deja affiche), cells[1:] = data rows
        for row in table.cells[1:]:
            table_lines.append("| " + " | ".join(row) + " |")
        parts.append("\n".join(table_lines))

    # Speaker notes
    if notes and len(notes) > 10:
        parts.append(f"\n--- Speaker Notes ---\n{notes}")

    return "\n\n".join(parts).strip()


def _is_section_slide(title: str, subtitle: str, content_text: str, notes: str) -> bool:
    """Detecte si la slide est un titre de section."""
    total = len(content_text.strip()) + len(notes.strip())
    return bool(title) and total < MIN_CONTENT_CHARS


def extract_pptx(
    file_path: str,
    tenant_id: str = "default",
    doc_id: Optional[str] = None,
) -> PptxExtractionResult:
    """
    Extrait un fichier PPTX et produit tous les artefacts du pipeline V2.

    Retourne VisionUnits + StructuralGraphBuildResult + full_text + structure,
    prets a etre injectes dans le pipeline d'extraction V2 a la place de Docling.

    Args:
        file_path: Chemin vers le fichier PPTX
        tenant_id: Tenant ID
        doc_id: Document ID (genere si absent)

    Returns:
        PptxExtractionResult avec tous les artefacts
    """
    from pptx import Presentation

    prs = Presentation(file_path)
    filename = Path(file_path).stem

    if not doc_id:
        doc_id = filename

    # Calculer le hash du document
    with open(file_path, "rb") as f:
        file_hash = hashlib.sha256(f.read()).hexdigest()
    doc_version_id = f"v1:{file_hash[:16]}"

    # Detecter le titre du document depuis la premiere slide
    doc_title = ""
    if prs.slides:
        first_title, first_subtitle = _extract_title_and_subtitle(prs.slides[0])
        if first_title and first_subtitle:
            doc_title = f"{first_title} {first_subtitle}".strip()
        elif first_title:
            doc_title = first_title
        else:
            doc_title = filename

    # =====================================================
    # Passe 1 : Extraire toutes les slides
    # =====================================================
    units: List[VisionUnit] = []
    doc_items: List[DocItem] = []
    chunks: List[TypeAwareChunk] = []
    sections: List[SectionInfo] = []
    page_outputs: List[PageOutput] = []
    page_indices: List[PageIndex] = []

    current_section = ""
    current_section_id = ""
    full_text_parts: List[str] = []
    full_text_offset = 0
    item_counter = 0

    for idx, slide in enumerate(prs.slides):
        title, subtitle = _extract_title_and_subtitle(slide)
        blocks, tables = _extract_content(slide, title, subtitle)
        notes = _extract_notes(slide)

        # Texte reconstruit de la slide
        reconstructed = _build_reconstructed_text(title, subtitle, blocks, tables, notes)

        # Detection role
        is_section = _is_section_slide(title, subtitle, reconstructed, notes)

        if is_section and idx > 0:
            current_section = title
            current_section_id = f"sec_{title.lower().replace(' ', '_')[:30]}_{uuid4().hex[:6]}"
            sections.append(SectionInfo(
                section_id=current_section_id,
                doc_id=doc_id or "",
                doc_version_id=doc_version_id,
                tenant_id=tenant_id,
                section_path=current_section,
                section_level=1,
                title=current_section,
            ))

        # Skip section slides sans contenu substantiel
        if is_section and len(notes) < MIN_CONTENT_CHARS and idx > 0:
            # Quand meme creer le VisionUnit pour coherence pipeline
            units.append(VisionUnit(
                id=f"PPTX_SLIDE_{idx}",
                format="PPTX",
                index=idx,
                dimensions=(DEFAULT_SLIDE_WIDTH, DEFAULT_SLIDE_HEIGHT),
                blocks=[],
                tables=[],
                notes="",
                has_notes=False,
                title=title,
            ))
            continue

        # Prefixe contextuel
        prefix_parts: List[str] = []
        if doc_title:
            prefix_parts.append(f"Document: {doc_title}")
        if current_section and current_section != title:
            prefix_parts.append(f"Section: {current_section}")
        prefix_parts.append(f"Slide {idx + 1}")

        prefixed_text = f"[{' | '.join(prefix_parts)}]\n\n{reconstructed}" if reconstructed else ""

        if len(prefixed_text.strip()) < 30:
            units.append(VisionUnit(
                id=f"PPTX_SLIDE_{idx}",
                format="PPTX",
                index=idx,
                dimensions=(DEFAULT_SLIDE_WIDTH, DEFAULT_SLIDE_HEIGHT),
                blocks=blocks,
                tables=tables,
                notes=notes,
                has_notes=bool(notes),
                title=title,
            ))
            continue

        # --- VisionUnit ---
        units.append(VisionUnit(
            id=f"PPTX_SLIDE_{idx}",
            format="PPTX",
            index=idx,
            dimensions=(DEFAULT_SLIDE_WIDTH, DEFAULT_SLIDE_HEIGHT),
            blocks=blocks,
            tables=tables,
            notes=notes,
            has_notes=bool(notes),
            title=title,
        ))

        # --- DocItems ---
        slide_item_ids: List[str] = []

        # Titre comme DocItem
        if title:
            item_id = f"item_{item_counter:04d}"
            item_counter += 1
            doc_items.append(DocItem(
                tenant_id=tenant_id,
                doc_id=doc_id,
                doc_version_id=doc_version_id,
                item_id=item_id,
                item_type=DocItemType.HEADING,
                text=title,
                page_no=idx,
                section_id=current_section_id or "",
                charspan_start=0,
                charspan_end=len(title),
                reading_order_index=item_counter,
            ))
            slide_item_ids.append(item_id)

        if subtitle:
            item_id = f"item_{item_counter:04d}"
            item_counter += 1
            doc_items.append(DocItem(
                tenant_id=tenant_id,
                doc_id=doc_id,
                doc_version_id=doc_version_id,
                item_id=item_id,
                item_type=DocItemType.HEADING,
                text=subtitle,
                page_no=idx,
                section_id=current_section_id or "",
                charspan_start=0,
                charspan_end=len(subtitle),
                reading_order_index=item_counter,
            ))
            slide_item_ids.append(item_id)

        # Contenu comme DocItems
        # Mapping TextBlock.type -> DocItemType
        _BLOCK_TYPE_MAP = {
            "list_item": DocItemType.LIST_ITEM,
            "paragraph": DocItemType.TEXT,
            "heading": DocItemType.HEADING,
            "table_cell": DocItemType.TABLE,
        }
        for block in blocks:
            item_id = f"item_{item_counter:04d}"
            item_counter += 1
            doc_items.append(DocItem(
                tenant_id=tenant_id,
                doc_id=doc_id,
                doc_version_id=doc_version_id,
                item_id=item_id,
                item_type=_BLOCK_TYPE_MAP.get(block.type, DocItemType.TEXT),
                text=block.text,
                page_no=idx,
                section_id=current_section_id or "",
                charspan_start=0,
                charspan_end=len(block.text),
                reading_order_index=item_counter,
            ))
            slide_item_ids.append(item_id)

        # Notes comme DocItem
        if notes and len(notes) > 10:
            item_id = f"item_{item_counter:04d}"
            item_counter += 1
            doc_items.append(DocItem(
                tenant_id=tenant_id,
                doc_id=doc_id,
                doc_version_id=doc_version_id,
                item_id=item_id,
                item_type=DocItemType.TEXT,
                text=notes,
                page_no=idx,
                section_id=current_section_id or "",
                charspan_start=0,
                charspan_end=len(notes),
                reading_order_index=item_counter,
            ))
            slide_item_ids.append(item_id)

        # --- TypeAwareChunk = slide reconstruite entiere ---
        # Chaque slide = 1 chunk autonome (la reconstruction preservee le contexte)
        chunk_kind = ChunkKind.NARRATIVE_TEXT
        if tables and not blocks:
            chunk_kind = ChunkKind.TABLE_TEXT

        chunks.append(TypeAwareChunk(
            chunk_id=f"chunk_slide_{idx:03d}_{uuid4().hex[:6]}",
            tenant_id=tenant_id,
            doc_id=doc_id,
            doc_version_id=doc_version_id,
            section_id=current_section_id or None,
            kind=chunk_kind,
            text=prefixed_text,
            item_ids=slide_item_ids,
            page_no=idx,
            page_span_min=idx,
            page_span_max=idx,
            is_relation_bearing=True,
            text_origin=TextOrigin.DOCLING,  # Semantiquement correct : extraction native
        ))

        # --- PageOutput pour la structure ---
        page_outputs.append(PageOutput(
            index=idx,
            text_markdown=prefixed_text,
        ))

        # --- full_text et PageIndex ---
        slide_text = f"\n[PAGE {idx + 1} | TYPE=slide]\n{prefixed_text}\n"
        page_indices.append(PageIndex(
            page_index=idx,
            start_offset=full_text_offset,
            end_offset=full_text_offset + len(slide_text),
            page_type="slide",
            title=title,
        ))
        full_text_parts.append(slide_text)
        full_text_offset += len(slide_text)

    # Assembler le full_text
    full_text = "\n".join(full_text_parts)

    # Construire le DocumentVersion
    doc_version = DocumentVersion(
        tenant_id=tenant_id,
        doc_id=doc_id,
        doc_version_id=doc_version_id,
        source_uri=file_path,
        title=doc_title,
        pipeline_version="v2_pptx_native",
        page_count=len(prs.slides),
        item_count=len(doc_items),
    )

    # Construire le StructuralGraphBuildResult
    structural_graph = StructuralGraphBuildResult(
        doc_items=doc_items,
        sections=sections,
        chunks=chunks,
        doc_version=doc_version,
        page_contexts=[],
        doc_dict={},  # Pas de DoclingDocument
    )

    # Structure
    structure = DocumentStructure(
        pages=page_outputs,
        metadata={
            "format": "PPTX",
            "extractor": "pptx_native",
            "slide_count": len(prs.slides),
            "doc_title": doc_title,
        },
    )

    logger.info(
        f"[OSMOSE:PptxExtractor] {len(prs.slides)} slides -> "
        f"{len(chunks)} chunks, {len(doc_items)} items "
        f"({sum(1 for u in units if u.has_notes)} with notes)"
    )

    return PptxExtractionResult(
        units=units,
        structural_graph=structural_graph,
        full_text=full_text,
        structure=structure,
        page_index=page_indices,
        doc_title=doc_title,
    )
