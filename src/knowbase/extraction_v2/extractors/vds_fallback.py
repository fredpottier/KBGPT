"""
VDSFallback - Fallback pour le signal VDS (Vector Drawing Signal).

Utilisé quand Docling ne fournit pas assez de détails sur les shapes/connecteurs.
- PDF: PyMuPDF page.get_drawings()
- PPTX: python-pptx MSO_SHAPE_TYPE

Spécification: OSMOSIS_EXTRACTION_V2_DECISIONS.md - Décision 6
"""

from __future__ import annotations
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple
import logging

logger = logging.getLogger(__name__)


class VDSFallback:
    """
    Fallback pour la détection VDS (Vector Drawing Signal).

    Utilisé quand Docling ne fournit pas assez d'informations
    sur les shapes vectoriels et connecteurs.

    Stratégie:
    - PDF: PyMuPDF page.get_drawings()
    - PPTX: python-pptx MSO_SHAPE_TYPE

    Usage:
        >>> fallback = VDSFallback()
        >>> info = fallback.analyze_pdf_page(pdf_path, page_index)
        >>> print(f"Connectors: {info['connector_count']}")
        >>> info = fallback.analyze_pptx_slide(pptx_path, slide_index)
        >>> print(f"Shapes: {info['shape_count']}")
    """

    def __init__(self):
        """Initialise le fallback VDS."""
        logger.info("[VDSFallback] Created")

    def analyze_pdf_page(
        self,
        pdf_path: str,
        page_index: int,
    ) -> Dict[str, Any]:
        """
        Analyse une page PDF pour les dessins vectoriels.

        Utilise PyMuPDF page.get_drawings() pour détecter:
        - Lignes
        - Rectangles
        - Courbes
        - Paths complexes

        Args:
            pdf_path: Chemin vers le PDF
            page_index: Index de la page (0-based)

        Returns:
            Dict avec:
            - connector_count: Nombre de connecteurs (lignes, flèches)
            - shape_count: Nombre de shapes (rectangles, etc.)
            - vector_density: Densité vectorielle (surface / page)
            - has_arrows: Présence de flèches détectées
        """
        try:
            import fitz  # PyMuPDF

            doc = fitz.open(pdf_path)
            if page_index >= len(doc):
                logger.warning(
                    f"[VDSFallback] Page {page_index} hors limites pour {pdf_path}"
                )
                return self._empty_result()

            page = doc[page_index]
            drawings = page.get_drawings()

            # Analyser les dessins
            connector_count = 0
            shape_count = 0
            total_area = 0.0
            has_arrows = False

            page_rect = page.rect
            page_area = page_rect.width * page_rect.height

            for drawing in drawings:
                # Type de dessin
                items = drawing.get("items", [])

                for item in items:
                    item_type = item[0] if item else None

                    if item_type == "l":  # Ligne
                        connector_count += 1
                        # Vérifier si c'est une flèche (heuristique simple)
                        if len(items) > 2:
                            has_arrows = True

                    elif item_type == "re":  # Rectangle
                        shape_count += 1
                        rect = item[1] if len(item) > 1 else None
                        if rect:
                            total_area += abs(rect.width * rect.height)

                    elif item_type == "c":  # Courbe
                        connector_count += 1

                    elif item_type == "qu":  # Quad
                        shape_count += 1

            doc.close()

            # Calculer la densité
            vector_density = min(total_area / page_area, 1.0) if page_area > 0 else 0.0

            result = {
                "connector_count": connector_count,
                "shape_count": shape_count,
                "vector_density": vector_density,
                "has_arrows": has_arrows,
                "total_drawings": len(drawings),
            }

            logger.debug(
                f"[VDSFallback] PDF page {page_index}: "
                f"{connector_count} connectors, {shape_count} shapes"
            )

            return result

        except ImportError:
            logger.warning("[VDSFallback] PyMuPDF non installé, utiliser: pip install PyMuPDF")
            return self._empty_result()
        except Exception as e:
            logger.error(f"[VDSFallback] Erreur analyse PDF: {e}")
            return self._empty_result()

    def analyze_pptx_slide(
        self,
        pptx_path: str,
        slide_index: int,
    ) -> Dict[str, Any]:
        """
        Analyse une slide PPTX pour les shapes vectoriels.

        Utilise python-pptx MSO_SHAPE_TYPE pour détecter:
        - Connecteurs (LINE, CONNECTOR)
        - Shapes (RECTANGLE, OVAL, etc.)
        - Groupes
        - Charts

        Args:
            pptx_path: Chemin vers le PPTX
            slide_index: Index de la slide (0-based)

        Returns:
            Dict avec:
            - connector_count: Nombre de connecteurs
            - shape_count: Nombre de shapes
            - group_count: Nombre de groupes
            - text_shape_count: Shapes avec texte
            - picture_count: Images
            - chart_count: Graphiques
        """
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE

            prs = Presentation(pptx_path)

            if slide_index >= len(prs.slides):
                logger.warning(
                    f"[VDSFallback] Slide {slide_index} hors limites pour {pptx_path}"
                )
                return self._empty_result_pptx()

            slide = prs.slides[slide_index]

            # Compteurs
            connector_count = 0
            shape_count = 0
            group_count = 0
            text_shape_count = 0
            picture_count = 0
            chart_count = 0

            for shape in slide.shapes:
                shape_type = shape.shape_type

                # Connecteurs (lignes, flèches)
                if shape_type in (
                    MSO_SHAPE_TYPE.LINE,
                    MSO_SHAPE_TYPE.FREEFORM,
                ):
                    connector_count += 1

                # Connecteurs spécifiques
                elif hasattr(MSO_SHAPE_TYPE, 'CONNECTOR'):
                    if shape_type == MSO_SHAPE_TYPE.CONNECTOR:
                        connector_count += 1

                # Groupes
                elif shape_type == MSO_SHAPE_TYPE.GROUP:
                    group_count += 1
                    # Compter les shapes dans le groupe
                    if hasattr(shape, 'shapes'):
                        for sub_shape in shape.shapes:
                            if sub_shape.shape_type in (
                                MSO_SHAPE_TYPE.LINE,
                                MSO_SHAPE_TYPE.FREEFORM,
                            ):
                                connector_count += 1
                            else:
                                shape_count += 1

                # Images
                elif shape_type == MSO_SHAPE_TYPE.PICTURE:
                    picture_count += 1

                # Graphiques
                elif shape_type == MSO_SHAPE_TYPE.CHART:
                    chart_count += 1

                # Autres shapes
                elif shape_type in (
                    MSO_SHAPE_TYPE.AUTO_SHAPE,
                    MSO_SHAPE_TYPE.TEXT_BOX,
                    MSO_SHAPE_TYPE.PLACEHOLDER,
                ):
                    shape_count += 1
                    # Vérifier si contient du texte
                    if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                        if shape.text_frame.text.strip():
                            text_shape_count += 1

                else:
                    shape_count += 1

            result = {
                "connector_count": connector_count,
                "shape_count": shape_count,
                "group_count": group_count,
                "text_shape_count": text_shape_count,
                "picture_count": picture_count,
                "chart_count": chart_count,
                "total_shapes": len(list(slide.shapes)),
            }

            logger.debug(
                f"[VDSFallback] PPTX slide {slide_index}: "
                f"{connector_count} connectors, {shape_count} shapes, "
                f"{group_count} groups"
            )

            return result

        except ImportError:
            logger.warning("[VDSFallback] python-pptx non installé")
            return self._empty_result_pptx()
        except Exception as e:
            logger.error(f"[VDSFallback] Erreur analyse PPTX: {e}")
            return self._empty_result_pptx()

    def count_connectors(
        self,
        file_path: str,
        page_or_slide_index: int,
    ) -> int:
        """
        Compte les connecteurs dans une page/slide.

        Détecte automatiquement le format (PDF ou PPTX).

        Args:
            file_path: Chemin vers le fichier
            page_or_slide_index: Index de la page/slide

        Returns:
            Nombre de connecteurs détectés
        """
        ext = Path(file_path).suffix.lower()

        if ext == ".pdf":
            info = self.analyze_pdf_page(file_path, page_or_slide_index)
        elif ext == ".pptx":
            info = self.analyze_pptx_slide(file_path, page_or_slide_index)
        else:
            logger.warning(f"[VDSFallback] Format non supporté: {ext}")
            return 0

        return info.get("connector_count", 0)

    def get_vector_density(
        self,
        file_path: str,
        page_or_slide_index: int,
    ) -> float:
        """
        Calcule la densité vectorielle d'une page/slide.

        Args:
            file_path: Chemin vers le fichier
            page_or_slide_index: Index de la page/slide

        Returns:
            Densité vectorielle (0.0 - 1.0)
        """
        ext = Path(file_path).suffix.lower()

        if ext == ".pdf":
            info = self.analyze_pdf_page(file_path, page_or_slide_index)
            return info.get("vector_density", 0.0)

        elif ext == ".pptx":
            info = self.analyze_pptx_slide(file_path, page_or_slide_index)
            # Estimer la densité basée sur le nombre de shapes
            total = info.get("total_shapes", 0)
            # Heuristique: 20+ shapes = densité élevée
            return min(total / 20.0, 1.0)

        return 0.0

    def has_diagram_indicators(
        self,
        file_path: str,
        page_or_slide_index: int,
    ) -> bool:
        """
        Vérifie si la page/slide a des indicateurs de diagramme.

        Critères:
        - Au moins 1 connecteur
        - OU au moins 5 shapes + 2 lignes

        Args:
            file_path: Chemin vers le fichier
            page_or_slide_index: Index de la page/slide

        Returns:
            True si indicateurs de diagramme détectés
        """
        ext = Path(file_path).suffix.lower()

        if ext == ".pdf":
            info = self.analyze_pdf_page(file_path, page_or_slide_index)
        elif ext == ".pptx":
            info = self.analyze_pptx_slide(file_path, page_or_slide_index)
        else:
            return False

        connectors = info.get("connector_count", 0)
        shapes = info.get("shape_count", 0)

        # Critères de détection
        if connectors >= 1:
            return True
        if shapes >= 5 and connectors >= 0:
            return True
        if info.get("has_arrows", False):
            return True

        return False

    @staticmethod
    def _empty_result() -> Dict[str, Any]:
        """Résultat vide pour PDF."""
        return {
            "connector_count": 0,
            "shape_count": 0,
            "vector_density": 0.0,
            "has_arrows": False,
            "total_drawings": 0,
        }

    @staticmethod
    def _empty_result_pptx() -> Dict[str, Any]:
        """Résultat vide pour PPTX."""
        return {
            "connector_count": 0,
            "shape_count": 0,
            "group_count": 0,
            "text_shape_count": 0,
            "picture_count": 0,
            "chart_count": 0,
            "total_shapes": 0,
        }


__all__ = ["VDSFallback"]
