"""
Extraction binaire du contenu PPTX.

Module extrait de pptx_pipeline.py pour réutilisabilité.
Supporte MegaParse (préféré) et python-pptx (fallback).
"""

import time
from pathlib import Path
from typing import List, Dict, Any, Optional
import logging


# Flags globaux pour détection des librairies disponibles
MEGAPARSE_AVAILABLE = False
PPTX_FALLBACK = False

try:
    from megaparse import MegaParse
    MEGAPARSE_AVAILABLE = True
except ImportError:
    pass

try:
    from pptx import Presentation
    PPTX_FALLBACK = True
except ImportError:
    pass


def extract_notes_and_text(
    pptx_path: Path,
    logger: Optional[logging.Logger] = None
) -> List[Dict[str, Any]]:
    """
    Point d'entrée principal pour extraction PPTX.
    Utilise MegaParse si disponible, sinon python-pptx.

    Args:
        pptx_path: Chemin vers le fichier PPTX
        logger: Logger optionnel

    Returns:
        List[Dict]: Liste des slides avec text, notes, metadata
    """
    if MEGAPARSE_AVAILABLE:
        return extract_with_megaparse(pptx_path, logger)
    elif PPTX_FALLBACK:
        return extract_with_python_pptx(pptx_path, logger)
    else:
        if logger:
            logger.error("❌ Aucune librairie d'extraction PPTX disponible!")
        return []


def extract_with_megaparse(
    pptx_path: Path,
    logger: Optional[logging.Logger] = None
) -> List[Dict[str, Any]]:
    """Extraction via MegaParse + python-pptx pour structure slides"""
    if logger:
        logger.info(f"📊 [MEGAPARSE] Extraction PPTX: {pptx_path.name}")

    try:
        # Étape 1: Obtenir la structure des slides via python-pptx
        if not PPTX_FALLBACK:
            raise ImportError("python-pptx requis pour obtenir la structure des slides")

        from pptx import Presentation
        prs = Presentation(str(pptx_path))
        slide_count = len(prs.slides)

        if logger:
            logger.info(f"📊 [MEGAPARSE] Détection {slide_count} slides via python-pptx")

        # Étape 2: Extraire le contenu enrichi via MegaParse
        megaparse = MegaParse()
        start_time = time.time()
        parsed_content = megaparse.load(str(pptx_path))
        load_duration = time.time() - start_time

        content_str = (
            str(parsed_content)
            if not isinstance(parsed_content, str)
            else parsed_content
        )

        # Étape 3: Diviser le contenu MegaParse en N slides
        slides_data = split_megaparse_by_slide_count(content_str, slide_count, pptx_path.name, logger)

        if logger:
            logger.info(
                f"✅ [MEGAPARSE] Extraction terminée - {len(slides_data)} slides en {load_duration:.1f}s"
            )
        return slides_data

    except Exception as e:
        if logger:
            logger.error(f"❌ [MEGAPARSE] Erreur: {str(e)}")

        if PPTX_FALLBACK:
            if logger:
                logger.info(f"🔄 Fallback vers python-pptx")
            return extract_with_python_pptx(pptx_path, logger)
        else:
            return [{
                "slide_index": 1,
                "text": f"Erreur d'extraction: {str(e)}",
                "notes": "",
                "megaparse_content": "",
                "content_type": "error",
            }]


def extract_with_python_pptx(
    pptx_path: Path,
    logger: Optional[logging.Logger] = None
) -> List[Dict[str, Any]]:
    """Extraction legacy via python-pptx"""
    if logger:
        logger.info(f"📊 Extraction via python-pptx: {pptx_path.name}")

    try:
        prs = Presentation(str(pptx_path))
        slides_data = []

        for i, slide in enumerate(prs.slides, start=1):
            notes = ""
            if getattr(slide, "has_notes_slide", False):
                notes_slide = getattr(slide, "notes_slide", None)
                if notes_slide and hasattr(notes_slide, "notes_text_frame"):
                    tf = notes_slide.notes_text_frame
                    if tf and hasattr(tf, "text"):
                        notes = (tf.text or "").strip()

            texts = []
            for shape in slide.shapes:
                txt = getattr(shape, "text", None)
                if isinstance(txt, str) and txt.strip():
                    texts.append(txt.strip())

                # Extraction tables
                if shape.has_table:
                    try:
                        table = shape.table
                        table_text = []
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                cell_text = (
                                    cell.text_frame.text.strip()
                                    if cell.text_frame
                                    else ""
                                )
                                if cell_text:
                                    row_text.append(cell_text)
                            if row_text:
                                table_text.append(" | ".join(row_text))
                        if table_text:
                            texts.append("[TABLE]\n" + "\n".join(table_text))
                    except Exception as e:
                        if logger:
                            logger.debug(f"Erreur extraction table slide {i}: {e}")

                # Extraction chart metadata
                if shape.shape_type == 3:  # MSO_SHAPE_TYPE.CHART
                    try:
                        chart_info = []
                        if hasattr(shape, "chart") and shape.chart:
                            chart = shape.chart
                            if hasattr(chart, "chart_title") and chart.chart_title:
                                title_text = (
                                    chart.chart_title.text_frame.text
                                    if hasattr(chart.chart_title, "text_frame")
                                    else ""
                                )
                                if title_text:
                                    chart_info.append(f"Chart Title: {title_text}")
                        if chart_info:
                            texts.append("[CHART]\n" + "\n".join(chart_info))
                    except Exception as e:
                        if logger:
                            logger.debug(f"Erreur extraction chart slide {i}: {e}")

            text_content = "\n".join(texts)
            slides_data.append({
                "slide_index": i,
                "text": text_content,
                "notes": notes,
                "megaparse_content": text_content,
                "content_type": "python_pptx_fallback",
            })

        if logger:
            logger.debug(f"Slides extraites via python-pptx: {len(slides_data)}")
        return slides_data

    except Exception as e:
        if logger:
            logger.error(f"❌ Erreur python-pptx: {str(e)}")
        return [{
            "slide_index": 1,
            "text": f"Erreur d'extraction python-pptx: {str(e)}",
            "notes": "",
            "megaparse_content": "",
            "content_type": "error",
        }]


def split_megaparse_by_slide_count(
    content: str,
    slide_count: int,
    source_name: str,
    logger: Optional[logging.Logger] = None
) -> List[Dict[str, Any]]:
    """
    Divise le contenu MegaParse en N slides basé sur le nombre réel de slides.

    LOGIQUE ORIGINALE (éprouvée) : Division proportionnelle par lignes.
    Le contenu MegaParse est divisé en N parties égales selon le nombre de lignes,
    où N = nombre de slides détecté par python-pptx.
    """
    slides_data = []

    # Diviser le contenu en lignes
    content_lines = content.split("\n")
    lines_per_slide = len(content_lines) // slide_count if slide_count > 0 else len(content_lines)

    if logger:
        logger.debug(f"[MEGAPARSE] Division {len(content_lines)} lignes en {slide_count} slides (~{lines_per_slide} lignes/slide)")

    for slide_num in range(1, slide_count + 1):
        # Calculer les indices de ligne pour cette slide
        start_line = (slide_num - 1) * lines_per_slide
        end_line = slide_num * lines_per_slide if slide_num < slide_count else len(content_lines)

        # Extraire le contenu pour cette slide
        slide_content = "\n".join(content_lines[start_line:end_line]).strip()

        # Ne créer une slide que si elle contient du contenu significatif
        if slide_content and len(slide_content) > 20:
            slides_data.append({
                "slide_index": slide_num,
                "text": slide_content,
                "notes": "",
                "megaparse_content": slide_content,
                "content_type": "megaparse_line_split",
            })

    if logger:
        logger.info(f"[MEGAPARSE] {len(slides_data)} slides créés pour {source_name}")

    return slides_data if slides_data else [{
        "slide_index": 1,
        "text": content,
        "notes": "",
        "megaparse_content": content,
        "content_type": "megaparse_single",
    }]


def extract_slides_from_megaparse(
    content: str,
    source_name: str,
    logger: Optional[logging.Logger] = None
) -> List[Dict[str, Any]]:
    """
    Legacy function - now redirects to split_megaparse_by_slide_count.
    Kept for backward compatibility.
    """
    # Essayer de détecter le nombre de slides depuis le contenu
    import re
    headers = list(re.finditer(r'^# .+$', content, re.MULTILINE))
    estimated_slide_count = max(1, len(headers))

    return split_megaparse_by_slide_count(content, estimated_slide_count, source_name, logger)
