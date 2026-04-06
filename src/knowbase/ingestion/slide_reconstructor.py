"""
Slide Reconstructor — Reconstitue un texte complet et autonome par slide PPTX.

Chaque slide est transformee en un "document texte" equivalent a une page PDF,
contenant toutes les couches d'information :
1. Contexte positionnel (titre de section, numero de slide)
2. Titre de la slide (+ sous-titre si present)
3. Contenu visible (bullets avec hierarchie d'indentation)
4. Tables (si presentes)
5. Notes orateur (verbatim auteur)

Le texte resultant doit passer le "test d'autonomie" : un expert humain
peut repondre a une question factuelle avec ce texte seul.

ADR: Unite de preuve vs Unite de lecture.
Les notes orateur sont du contenu auteur verbatim — zero tension avec l'invariant.
"""
from __future__ import annotations

import logging
import re
from dataclasses import dataclass
from enum import Enum

logger = logging.getLogger(__name__)

# Seuil minimum pour considerer une slide comme "contenu"
# (vs slide de titre de section)
MIN_CONTENT_CHARS = 50


class SlideRole(Enum):
    """Role detecte d'une slide dans la presentation."""
    SECTION_TITLE = "section_title"
    CONTENT = "content"


@dataclass
class ReconstructedSlide:
    """Texte reconstruit d'une slide PPTX."""
    slide_index: int
    title: str
    subtitle: str
    text: str
    has_notes: bool
    notes_length: int
    content_length: int
    section_title: str = ""
    role: SlideRole = SlideRole.CONTENT


def _extract_title_and_subtitle(slide) -> tuple[str, str]:
    """Extrait titre et sous-titre en gerant les placeholders PPTX.

    PPTX utilise \\x0b (vertical tab) comme saut de ligne dans les titres.
    On le traite comme separateur titre/sous-titre.
    """
    title = ""
    subtitle = ""

    for shape in slide.placeholders:
        phf = shape.placeholder_format
        if phf is None:
            continue
        # idx 0 = titre, idx 1 = sous-titre/corps
        if phf.idx == 0 and shape.has_text_frame:
            raw = shape.text_frame.text.strip()
            # \x0b = vertical tab, utilise par PPTX pour les sauts de ligne dans les titres
            if "\x0b" in raw:
                parts = raw.split("\x0b", 1)
                title = parts[0].strip()
                subtitle = parts[1].strip()
            else:
                title = raw
        elif phf.idx == 1 and shape.has_text_frame:
            candidate = shape.text_frame.text.strip()
            if candidate and len(candidate) < 200 and candidate.count("\n") <= 2:
                if not subtitle:
                    subtitle = candidate

    # Fallback sur shapes.title si pas de placeholder 0
    if not title and slide.shapes.title:
        raw = slide.shapes.title.text.strip()
        if "\x0b" in raw:
            parts = raw.split("\x0b", 1)
            title = parts[0].strip()
            if not subtitle:
                subtitle = parts[1].strip()
        else:
            title = raw

    return title, subtitle


def _detect_slide_role(title: str, subtitle: str, content_text: str, notes_text: str) -> SlideRole:
    """Detecte si la slide est un titre de section ou du contenu."""
    total_content = len(content_text.strip()) + len(notes_text.strip())
    # Slide de section : titre present, peu ou pas de contenu au-dela
    if title and total_content < MIN_CONTENT_CHARS:
        return SlideRole.SECTION_TITLE
    return SlideRole.CONTENT


def _clean_whitespace(text: str) -> str:
    """Normalise les espaces multiples et tabs dans le texte extrait."""
    text = text.replace("\t", " ")
    text = re.sub(r" {2,}", " ", text)
    return text.strip()


def _extract_content_lines(slide, title: str, subtitle: str) -> list[str]:
    """Extrait le contenu visible (bullets + tables) en preservant la hierarchie."""
    skip_texts = {title, subtitle}
    content_lines: list[str] = []

    # Identifier les placeholders titre/sous-titre a ignorer
    title_shape_ids: set[int] = set()
    for shape in slide.placeholders:
        phf = shape.placeholder_format
        if phf is not None and phf.idx in (0, 1):
            title_shape_ids.add(shape.shape_id)

    for shape in slide.shapes:
        # Ignorer les shapes titre/sous-titre (deja extraits)
        if shape.shape_id in title_shape_ids:
            continue

        if shape.has_text_frame:
            shape_lines: list[str] = []
            for para in shape.text_frame.paragraphs:
                text = _clean_whitespace(para.text)
                if not text or text in skip_texts:
                    continue
                level = para.level
                if level > 0:
                    prefix = "  " * level + "- "
                else:
                    prefix = "- " if len(text) < 300 else ""
                shape_lines.append(f"{prefix}{text}")
            if shape_lines:
                # Filtrer les shapes qui ne contiennent que des labels courts
                # (typiquement des elements de diagramme : "KPI 1", "Process", etc.)
                avg_len = sum(len(l.lstrip(" -")) for l in shape_lines) / len(shape_lines)
                if avg_len < 15 and len(shape_lines) > 2:
                    continue
                content_lines.extend(shape_lines)

        if shape.has_table:
            table = shape.table
            rows = list(table.rows)
            if not rows:
                continue
            table_lines: list[str] = []
            headers = [_clean_whitespace(cell.text) for cell in rows[0].cells]
            table_lines.append("| " + " | ".join(headers) + " |")
            table_lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
            for row in rows[1:]:
                cells = [_clean_whitespace(cell.text) for cell in row.cells]
                table_lines.append("| " + " | ".join(cells) + " |")
            content_lines.append("\n".join(table_lines))

    return content_lines


def _extract_notes(slide) -> str:
    """Extrait les notes orateur."""
    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
        return slide.notes_slide.notes_text_frame.text.strip()
    return ""


def reconstruct_slides_from_pptx(pptx_path: str, doc_title: str = "") -> list[ReconstructedSlide]:
    """
    Reconstruit un texte complet et autonome pour chaque slide d'un PPTX.

    Args:
        pptx_path: Chemin vers le fichier PPTX
        doc_title: Titre du document (pour le prefixe contextuel)

    Returns:
        Liste de ReconstructedSlide, une par slide non-vide
    """
    from pptx import Presentation

    prs = Presentation(pptx_path)
    results: list[ReconstructedSlide] = []
    current_section = ""

    for idx, slide in enumerate(prs.slides):
        title, subtitle = _extract_title_and_subtitle(slide)
        content_lines = _extract_content_lines(slide, title, subtitle)
        notes_text = _extract_notes(slide)

        # Detecter le role de la slide
        content_text = "\n".join(content_lines)
        role = _detect_slide_role(title, subtitle, content_text, notes_text)

        # Mettre a jour la section courante
        if role == SlideRole.SECTION_TITLE and idx > 0:
            # idx > 0 : la slide de couverture n'est pas une section
            current_section = title
            # Les slides de section ne sont pas emises comme chunks autonomes
            # sauf si elles ont des notes substantielles
            if len(notes_text) < MIN_CONTENT_CHARS:
                continue

        # Construire le texte reconstruit
        parts: list[str] = []

        # Titre (+ sous-titre separe)
        if title:
            header = f"## {title}"
            if subtitle:
                header += f"\n### {subtitle}"
            parts.append(header)
        elif subtitle:
            parts.append(f"### {subtitle}")

        if content_lines:
            parts.append("\n".join(content_lines))

        if notes_text and len(notes_text) > 10:
            parts.append(f"\n--- Speaker Notes ---\n{notes_text}")

        full_text = "\n\n".join(parts).strip()

        # Filtrer les slides sans contenu exploitable
        if len(full_text) < 30:
            continue

        # Prefixe contextuel
        prefix_parts: list[str] = []
        if doc_title:
            prefix_parts.append(f"Document: {doc_title}")
        if current_section and current_section != title:
            prefix_parts.append(f"Section: {current_section}")
        prefix_parts.append(f"Slide {idx + 1}")

        prefixed_text = f"[{' | '.join(prefix_parts)}]\n\n{full_text}"

        results.append(ReconstructedSlide(
            slide_index=idx,
            title=title,
            subtitle=subtitle,
            text=prefixed_text,
            has_notes=bool(notes_text),
            notes_length=len(notes_text),
            content_length=len(full_text),
            section_title=current_section,
            role=role,
        ))

    logger.info(
        f"[SlideReconstructor] {len(prs.slides)} slides -> {len(results)} reconstructed "
        f"({sum(1 for s in results if s.has_notes)} with notes)"
    )

    return results
