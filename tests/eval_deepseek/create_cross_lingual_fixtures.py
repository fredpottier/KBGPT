"""
Script Helper: Création automatique des fixtures cross-lingual pour Phase 3

Génère slides PPTX simples en EN/FR/DE pour tester la préservation USP cross-lingual
"""

from pathlib import Path
from typing import Dict, List
import json

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("⚠️ python-pptx non installé. Installer avec: pip install python-pptx")


class CrossLingualFixtureGenerator:
    """Générateur de fixtures cross-lingual pour tests OSMOSE"""

    def __init__(self, output_dir: Path = None):
        if output_dir is None:
            output_dir = Path("tests/eval_deepseek/fixtures/cross_lingual")
        self.output_dir = output_dir
        self.output_dir.mkdir(parents=True, exist_ok=True)

    def create_simple_slide(
        self,
        title: str,
        content: str,
        output_path: Path,
        background_color: tuple = (255, 255, 255)
    ):
        """
        Créer slide PPTX simple avec titre + contenu

        Args:
            title: Titre slide
            content: Contenu principal
            output_path: Path de sauvegarde
            background_color: RGB tuple (default: blanc)
        """
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx requis")

        # Créer présentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # Blank slide layout
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*background_color)

        # Title box
        title_box = slide.shapes.add_textbox(
            left=Inches(0.5),
            top=Inches(0.5),
            width=Inches(9),
            height=Inches(1.2)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
        title_para.alignment = PP_ALIGN.CENTER

        # Content box
        content_box = slide.shapes.add_textbox(
            left=Inches(0.5),
            top=Inches(2.5),
            width=Inches(9),
            height=Inches(4)
        )
        content_frame = content_box.text_frame
        content_frame.text = content
        content_frame.word_wrap = True
        content_para = content_frame.paragraphs[0]
        content_para.font.size = Pt(20)
        content_para.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray
        content_para.alignment = PP_ALIGN.LEFT
        content_para.space_before = Pt(12)
        content_para.line_spacing = 1.5

        # Save
        prs.save(str(output_path))
        print(f"   ✅ Créé: {output_path.name}")

    def generate_crr_fixtures(self) -> Dict[str, Path]:
        """
        Générer fixtures pour concept 'Customer Retention Rate'

        Returns:
            Dict avec paths EN/FR/DE
        """
        print("\n📄 Génération fixtures: Customer Retention Rate (CRR)")

        fixtures = {}

        # English
        fixtures["en"] = self.output_dir / "crr_definition_en.pptx"
        self.create_simple_slide(
            title="Customer Retention Rate (CRR)",
            content=(
                "Definition:\n"
                "The Customer Retention Rate (CRR) measures the percentage of customers "
                "retained over a specific period.\n\n"
                "Formula: CRR = ((E - N) / S) × 100\n\n"
                "Where:\n"
                "• E = customers at end of period\n"
                "• N = new customers acquired during period\n"
                "• S = customers at start of period\n\n"
                "Example:\n"
                "A company starts with 500 customers (S), gains 100 new customers (N), "
                "and ends with 550 customers (E).\n"
                "CRR = ((550 - 100) / 500) × 100 = 90%"
            ),
            output_path=fixtures["en"]
        )

        # French
        fixtures["fr"] = self.output_dir / "crr_definition_fr.pptx"
        self.create_simple_slide(
            title="Taux de Rétention Client (CRR)",
            content=(
                "Définition:\n"
                "Le Taux de Rétention Client (CRR) mesure le pourcentage de clients "
                "conservés sur une période donnée.\n\n"
                "Formule: CRR = ((E - N) / S) × 100\n\n"
                "Où:\n"
                "• E = clients en fin de période\n"
                "• N = nouveaux clients acquis pendant la période\n"
                "• S = clients au début de période\n\n"
                "Exemple:\n"
                "Une entreprise commence avec 500 clients (S), acquiert 100 nouveaux "
                "clients (N), et termine avec 550 clients (E).\n"
                "CRR = ((550 - 100) / 500) × 100 = 90%"
            ),
            output_path=fixtures["fr"]
        )

        # German
        fixtures["de"] = self.output_dir / "crr_definition_de.pptx"
        self.create_simple_slide(
            title="Kundenbindungsrate (CRR)",
            content=(
                "Definition:\n"
                "Die Kundenbindungsrate (CRR) misst den Prozentsatz der Kunden, "
                "die über einen bestimmten Zeitraum gehalten werden.\n\n"
                "Formel: CRR = ((E - N) / S) × 100\n\n"
                "Wobei:\n"
                "• E = Kunden am Ende des Zeitraums\n"
                "• N = neue Kunden im Zeitraum gewonnen\n"
                "• S = Kunden zu Beginn des Zeitraums\n\n"
                "Beispiel:\n"
                "Ein Unternehmen beginnt mit 500 Kunden (S), gewinnt 100 neue "
                "Kunden (N), und endet mit 550 Kunden (E).\n"
                "CRR = ((550 - 100) / 500) × 100 = 90%"
            ),
            output_path=fixtures["de"]
        )

        return fixtures

    def generate_auth_policy_fixtures(self) -> Dict[str, Path]:
        """
        Générer fixtures pour concept 'Authentication Policy'

        Returns:
            Dict avec paths EN/FR/DE
        """
        print("\n📄 Génération fixtures: Authentication Policy")

        fixtures = {}

        # English
        fixtures["en"] = self.output_dir / "auth_policy_en.pptx"
        self.create_simple_slide(
            title="Multi-Factor Authentication Policy",
            content=(
                "Security Requirement:\n"
                "All users accessing sensitive systems MUST use Multi-Factor "
                "Authentication (MFA).\n\n"
                "Implementation:\n"
                "• Primary factor: Password (minimum 12 characters)\n"
                "• Secondary factor: One of the following:\n"
                "  - SMS verification code\n"
                "  - Authenticator app (TOTP)\n"
                "  - Hardware security key (U2F)\n\n"
                "Compliance:\n"
                "This policy aligns with ISO 27001 Access Control requirements "
                "and NIST 800-63B guidelines."
            ),
            output_path=fixtures["en"]
        )

        # French
        fixtures["fr"] = self.output_dir / "auth_policy_fr.pptx"
        self.create_simple_slide(
            title="Politique d'Authentification Multi-Facteurs",
            content=(
                "Exigence de Sécurité:\n"
                "Tous les utilisateurs accédant aux systèmes sensibles DOIVENT utiliser "
                "l'Authentification Multi-Facteurs (AMF).\n\n"
                "Implémentation:\n"
                "• Facteur primaire: Mot de passe (minimum 12 caractères)\n"
                "• Facteur secondaire: Un des suivants:\n"
                "  - Code de vérification SMS\n"
                "  - Application d'authentification (TOTP)\n"
                "  - Clé de sécurité matérielle (U2F)\n\n"
                "Conformité:\n"
                "Cette politique s'aligne avec les exigences de Contrôle d'Accès "
                "ISO 27001 et les directives NIST 800-63B."
            ),
            output_path=fixtures["fr"]
        )

        # German
        fixtures["de"] = self.output_dir / "auth_policy_de.pptx"
        self.create_simple_slide(
            title="Richtlinie zur Multi-Faktor-Authentifizierung",
            content=(
                "Sicherheitsanforderung:\n"
                "Alle Benutzer, die auf sensible Systeme zugreifen, MÜSSEN "
                "Multi-Faktor-Authentifizierung (MFA) verwenden.\n\n"
                "Implementierung:\n"
                "• Primärer Faktor: Passwort (mindestens 12 Zeichen)\n"
                "• Sekundärer Faktor: Einer der folgenden:\n"
                "  - SMS-Bestätigungscode\n"
                "  - Authentifizierungs-App (TOTP)\n"
                "  - Hardware-Sicherheitsschlüssel (U2F)\n\n"
                "Compliance:\n"
                "Diese Richtlinie entspricht den ISO 27001 Zugangskontrollanforderungen "
                "und NIST 800-63B Richtlinien."
            ),
            output_path=fixtures["de"]
        )

        return fixtures

    def generate_all_fixtures(self) -> Dict[str, Dict[str, Path]]:
        """
        Générer tous les fixtures nécessaires pour Phase 3

        Returns:
            Dict avec tous les fixtures générés
        """
        print("\n" + "="*60)
        print("🎨 Génération Fixtures Cross-Lingual pour Phase 3")
        print("="*60)

        if not PPTX_AVAILABLE:
            print("\n❌ python-pptx non installé")
            print("   Installer avec: pip install python-pptx")
            return {}

        all_fixtures = {}

        # CRR fixtures
        all_fixtures["crr"] = self.generate_crr_fixtures()

        # Auth Policy fixtures
        all_fixtures["auth_policy"] = self.generate_auth_policy_fixtures()

        # Summary
        print("\n" + "="*60)
        print("✅ Génération Terminée")
        print("="*60)
        print(f"\nLocation: {self.output_dir}")
        print(f"\nFichiers créés:")
        for concept, fixtures in all_fixtures.items():
            print(f"\n{concept}:")
            for lang, path in fixtures.items():
                print(f"  • {lang}: {path.name}")

        # Sauvegarder manifest
        manifest = {
            "generated_at": __import__("time").strftime("%Y-%m-%d %H:%M:%S"),
            "output_dir": str(self.output_dir),
            "fixtures": {
                concept: {lang: str(path) for lang, path in fixtures.items()}
                for concept, fixtures in all_fixtures.items()
            }
        }

        manifest_path = self.output_dir / "fixtures_manifest.json"
        with open(manifest_path, 'w', encoding='utf-8') as f:
            json.dump(manifest, f, indent=2, ensure_ascii=False)
        print(f"\n💾 Manifest: {manifest_path}")

        # Instructions next steps
        print("\n📝 Prochaines Étapes:")
        print("   1. Vérifier les PPTX générés (ouvrir dans PowerPoint)")
        print("   2. Convertir PPTX → PNG si nécessaire:")
        print("      libreoffice --headless --convert-to png *.pptx")
        print("   3. Lancer Phase 3:")
        print("      cd tests/eval_deepseek")
        print("      python test_03_cross_lingual.py")

        return all_fixtures

    def convert_pptx_to_png(self, pptx_path: Path) -> Path:
        """
        Convertir PPTX → PNG (première slide)

        Note: Nécessite LibreOffice ou pdf2image
        TODO: Implémenter conversion automatique
        """
        print(f"⚠️ Conversion PPTX → PNG à implémenter")
        print(f"   Manuel: Ouvrir {pptx_path} et exporter en PNG")
        print(f"   Ou: libreoffice --headless --convert-to png {pptx_path}")
        return pptx_path.with_suffix(".png")


def main():
    """Point d'entrée principal"""
    generator = CrossLingualFixtureGenerator()

    # Generate all
    fixtures = generator.generate_all_fixtures()

    if fixtures:
        print("\n✅ SUCCESS - Fixtures prêtes pour Phase 3")
        return 0
    else:
        print("\n❌ FAIL - Vérifier installation python-pptx")
        return 1


if __name__ == "__main__":
    exit(main())
